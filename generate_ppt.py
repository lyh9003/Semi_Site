"""
SemiCon 프로젝트 구조 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── 색상 팔레트 ───────────────────────────────────────────────────────────────
C_BG_DARK    = RGBColor(0x1e, 0x29, 0x3b)   # slate-800
C_BG_DARKER  = RGBColor(0x0f, 0x17, 0x2a)   # slate-900
C_BLUE       = RGBColor(0x37, 0x7d, 0xf7)   # blue-500
C_BLUE_DARK  = RGBColor(0x1d, 0x4e, 0xd8)   # blue-700
C_PURPLE     = RGBColor(0x7c, 0x3a, 0xed)   # purple-600
C_TEAL       = RGBColor(0x0d, 0x94, 0x88)   # teal-600
C_GREEN      = RGBColor(0x10, 0xb9, 0x81)   # emerald-500
C_ORANGE     = RGBColor(0xf5, 0x9e, 0x0b)   # amber-500
C_RED        = RGBColor(0xef, 0x44, 0x44)   # red-500
C_WHITE      = RGBColor(0xff, 0xff, 0xff)
C_SLATE_200  = RGBColor(0xe2, 0xe8, 0xf0)
C_SLATE_400  = RGBColor(0x94, 0xa3, 0xb8)
C_SLATE_600  = RGBColor(0x47, 0x55, 0x69)
C_SLATE_700  = RGBColor(0x33, 0x41, 0x55)
C_ACCENT     = RGBColor(0x60, 0xa5, 0xfa)   # blue-400


# ─── 헬퍼 함수 ────────────────────────────────────────────────────────────────
def blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, l, t, w, h, fill=None, line=None, radius=Pt(8)):
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.adjustments[0] = 0.05
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return txBox

def add_arrow(slide, x1, y1, x2, y2, color=C_SLATE_400, w=Pt(1.5)):
    """수평/수직 화살표 (connector)"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = w
    return connector

def section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_BG_DARKER)
    # 좌측 강조 바
    add_rect(slide, 0, 0, 0.12, 7.5, fill=C_BLUE)
    add_text(slide, title, 0.5, 2.8, 12, 1.0, size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.5, 3.9, 12, 0.6, size=18, color=C_SLATE_400, align=PP_ALIGN.CENTER)

def slide_header(slide, title, accent_color=C_BLUE):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=C_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=C_BG_DARKER)
    add_rect(slide, 0, 0, 0.08, 1.0, fill=accent_color)
    add_text(slide, title, 0.25, 0.18, 12, 0.65, size=22, bold=True, color=C_WHITE)

def badge(slide, text, l, t, w, h, bg=C_BLUE, text_color=C_WHITE, size=9):
    add_rounded_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, text, l, t, w, h, size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def info_box(slide, title, lines, l, t, w, h, accent=C_BLUE, title_size=11, body_size=9):
    add_rounded_rect(slide, l, t, w, h, fill=C_SLATE_700, line=accent)
    add_text(slide, title, l+0.1, t+0.08, w-0.2, 0.28, size=title_size, bold=True, color=accent)
    body = "\n".join(lines)
    add_text(slide, body, l+0.1, t+0.38, w-0.2, h-0.48, size=body_size, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_BG_DARKER)
# 그라디언트 효과 (상단 줄)
add_rect(slide, 0, 0, 13.33, 0.08, fill=C_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill=C_BLUE)

# 중앙 글로우 효과 배경
add_rounded_rect(slide, 2.5, 1.8, 8.33, 4.5, fill=C_BG_DARK, line=C_BLUE_DARK)

add_text(slide, "⚡ SemiCon", 2.5, 2.0, 8.33, 0.8, size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "반도체 시황 인텔리전스 플랫폼", 2.5, 2.65, 8.33, 1.1, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "시스템 아키텍처 & 기술 구조 분석", 2.5, 3.6, 8.33, 0.6, size=18, color=C_SLATE_400, align=PP_ALIGN.CENTER)

# 특징 배지
badges = [("RAG Q&A", C_BLUE), ("AI 브리핑", C_PURPLE), ("지식 그래프", C_TEAL), ("AI 에이전트", C_GREEN)]
bx = 2.8
for label, col in badges:
    add_rounded_rect(slide, bx, 4.5, 1.7, 0.42, fill=col)
    add_text(slide, label, bx, 4.5, 1.7, 0.42, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    bx += 1.95

add_text(slide, "Next.js 15  ·  Supabase  ·  OpenAI  ·  D3.js  ·  FastAPI  ·  Ollama",
         2.5, 5.3, 8.33, 0.4, size=11, color=C_SLATE_400, align=PP_ALIGN.CENTER)
add_text(slide, "2026.05", 0, 7.0, 13.33, 0.4, size=10, color=C_SLATE_600, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "목차 — 발표 순서")

items = [
    ("01", "프로젝트 개요 & 전체 아키텍처", C_BLUE),
    ("02", "기술 스택 & 인프라 구조", C_PURPLE),
    ("03", "페이지 & 라우트 구조", C_TEAL),
    ("04", "RAG Q&A 챗봇 시스템 상세", C_BLUE),
    ("05", "일일 브리핑 시스템", C_ORANGE),
    ("06", "AI 에이전트 시황 채팅", C_GREEN),
    ("07", "지식 그래프 & 인사이트 시스템", C_PURPLE),
    ("08", "Supabase DB 구조 & SQL 함수", C_TEAL),
    ("09", "데이터 현황 & 비용 구조", C_GREEN),
]

for i, (num, title, col) in enumerate(items):
    row = i % 5
    col_idx = i // 5
    lx = 0.3 + col_idx * 6.5
    ty = 1.15 + row * 1.2
    add_rounded_rect(slide, lx, ty, 6.1, 1.0, fill=C_SLATE_700, line=col)
    add_rect(slide, lx, ty, 0.7, 1.0, fill=col)
    add_text(slide, num, lx, ty, 0.7, 1.0, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+0.8, ty+0.28, 5.1, 0.5, size=13, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 섹션 헤더: 개요 & 아키텍처
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "01", "프로젝트 개요 & 전체 아키텍처")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 프로젝트 개요
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "프로젝트 개요 — SemiCon이란?")

# 좌측: 미션
add_rounded_rect(slide, 0.25, 1.2, 5.8, 5.9, fill=C_SLATE_700, line=C_BLUE)
add_text(slide, "🎯  미션", 0.45, 1.35, 5.4, 0.4, size=14, bold=True, color=C_ACCENT)
mission_lines = [
    "반도체 시황 정보의 파편화 문제를 해결한다.",
    "",
    "뉴스 · 증권리포트 · 텔레그램 채널에 흩어진",
    "정보를 하나의 플랫폼으로 통합하고,",
    "AI가 즉각적인 시황 인사이트를 제공한다.",
]
add_text(slide, "\n".join(mission_lines), 0.45, 1.85, 5.4, 2.0, size=12, color=C_SLATE_200)

add_text(slide, "📊  주요 데이터 규모", 0.45, 3.9, 5.4, 0.4, size=13, bold=True, color=C_ACCENT)
stats = [
    ("뉴스 기사", "~30,000건", C_BLUE),
    ("증권 리포트", "~584건", C_PURPLE),
    ("텔레그램 메시지", "~5,558건", C_TEAL),
    ("지식 그래프 엔티티", "2,121개", C_GREEN),
]
for j, (label, val, col) in enumerate(stats):
    ty = 4.45 + j * 0.55
    add_rect(slide, 0.45, ty, 0.06, 0.38, fill=col)
    add_text(slide, label, 0.6, ty, 2.5, 0.38, size=11, color=C_SLATE_400)
    add_text(slide, val, 3.5, ty, 2.2, 0.38, size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)

# 우측: 핵심 기능
add_rounded_rect(slide, 6.4, 1.2, 6.65, 5.9, fill=C_SLATE_700, line=C_PURPLE)
add_text(slide, "🚀  핵심 기능 6가지", 6.6, 1.35, 6.2, 0.4, size=14, bold=True, color=RGBColor(0xc4, 0xb5, 0xfd))

features = [
    ("📰 뉴스 & 리포트", "최신 반도체 뉴스·증권리포트 통합 열람", C_BLUE),
    ("💬 RAG Q&A", "자연어 질문 → AI 답변 + 출처 카드", C_BLUE),
    ("📊 일일 브리핑", "매일 최신 시황을 AI가 자동 요약", C_ORANGE),
    ("🕸️ 지식 그래프", "엔티티 관계 D3.js 인터랙티브 시각화", C_PURPLE),
    ("🔗 인사이트", "키워드·의미 유사도 교차 참조", C_TEAL),
    ("🤖 AI 에이전트", "5개 에이전트가 실시간 시황 토론", C_GREEN),
]
for j, (title, desc, col) in enumerate(features):
    ty = 1.95 + j * 0.87
    add_rounded_rect(slide, 6.55, ty, 6.3, 0.78, fill=C_BG_DARK, line=col)
    add_text(slide, title, 6.72, ty+0.08, 2.5, 0.35, size=11, bold=True, color=col)
    add_text(slide, desc, 6.72, ty+0.4, 5.9, 0.32, size=9.5, color=C_SLATE_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 전체 아키텍처 다이어그램
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "전체 시스템 아키텍처")

add_rect(slide, 0, 0, 13.33, 7.5, fill=C_BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.0, fill=C_BG_DARKER)
add_rect(slide, 0, 0, 0.08, 1.0, fill=C_BLUE)
add_text(slide, "전체 시스템 아키텍처", 0.25, 0.18, 12, 0.65, size=22, bold=True, color=C_WHITE)

# ── 레이어 1: 사용자 ──
add_text(slide, "👤  사용자 (브라우저)", 0.3, 1.05, 3, 0.3, size=10, bold=True, color=C_SLATE_400)
add_rounded_rect(slide, 0.3, 1.35, 3.2, 0.55, fill=C_BG_DARKER, line=C_SLATE_600)
add_text(slide, "브라우저 / 웹 클라이언트", 0.3, 1.35, 3.2, 0.55, size=10, color=C_SLATE_200, align=PP_ALIGN.CENTER)

# ── 레이어 2: Next.js (Vercel) ──
add_text(slide, "🌐  Next.js 15 App (Vercel)", 0.3, 2.1, 8, 0.3, size=10, bold=True, color=C_ACCENT)
add_rounded_rect(slide, 0.3, 2.38, 12.7, 1.5, fill=C_SLATE_700, line=C_BLUE)

pages = [
    ("/ 홈\n브리핑", 0.45, 2.52),
    ("/news\n뉴스", 1.85, 2.52),
    ("/reports\n리포트", 3.25, 2.52),
    ("/ask\nQ&A", 4.65, 2.52),
    ("/graph\n지식그래프", 6.05, 2.52),
    ("/insight\n인사이트", 7.55, 2.52),
    ("/agents\nAI채팅", 9.05, 2.52),
    ("/board\n게시판", 10.55, 2.52),
]
for label, lx, ty in pages:
    add_rounded_rect(slide, lx, ty, 1.25, 1.1, fill=C_BG_DARK, line=C_SLATE_600)
    add_text(slide, label, lx, ty+0.05, 1.25, 1.0, size=8.5, color=C_SLATE_200, align=PP_ALIGN.CENTER)

# ── API 레이어 ──
add_text(slide, "⚙️  API Routes (Next.js)", 0.3, 4.05, 8, 0.3, size=10, bold=True, color=C_ACCENT)
add_rounded_rect(slide, 0.3, 4.32, 12.7, 0.9, fill=C_SLATE_700, line=C_BLUE_DARK)

apis = [
    ("/api/chat\nRAG Q&A", C_BLUE),
    ("/api/briefing\n일일브리핑", C_ORANGE),
    ("/api/graph/summary\nAI요약", C_PURPLE),
    ("/api/insight/similar\n의미유사도", C_TEAL),
    ("/api/graph\n그래프데이터", C_PURPLE),
    ("/api/board\n게시판CRUD", C_SLATE_600),
    ("/api/stocks\n주가데이터", C_GREEN),
]
ax = 0.45
for label, col in apis:
    add_rounded_rect(slide, ax, 4.45, 1.7, 0.68, fill=C_BG_DARK, line=col)
    add_text(slide, label, ax, 4.48, 1.7, 0.62, size=7.5, color=C_SLATE_200, align=PP_ALIGN.CENTER)
    ax += 1.82

# ── 외부 서비스 레이어 ──
add_text(slide, "🔌  외부 서비스 & 데이터 소스", 0.3, 5.35, 8, 0.3, size=10, bold=True, color=C_SLATE_400)

services = [
    ("🗄️ Supabase\nPostgreSQL+pgvector", 0.3, 5.62, 2.4, C_TEAL),
    ("🤖 OpenAI API\nGPT-4o-mini + Embeddings", 2.85, 5.62, 2.4, C_GREEN),
    ("🦙 Ollama (로컬)\nqwen2.5:7b", 5.4, 5.62, 2.0, C_ORANGE),
    ("🐍 FastAPI WS\nagents_server.py", 7.55, 5.62, 2.2, C_PURPLE),
    ("⚡ Vercel CDN\n엣지캐시", 9.9, 5.62, 1.7, C_BLUE),
    ("📱 카카오 OAuth\n소셜로그인", 11.75, 5.62, 1.5, C_YELLOW if False else C_ORANGE),
]
for label, lx, ty, w, col in services:
    add_rounded_rect(slide, lx, ty, w, 1.0, fill=C_BG_DARKER, line=col)
    add_text(slide, label, lx, ty+0.08, w, 0.85, size=8.5, color=C_SLATE_200, align=PP_ALIGN.CENTER)

# 화살표
add_arrow(slide, 1.9, 1.9, 1.9, 2.38, color=C_BLUE)
add_arrow(slide, 6.65, 3.88, 6.65, 4.32, color=C_BLUE)
add_arrow(slide, 6.65, 5.22, 6.65, 5.62, color=C_SLATE_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 섹션 헤더: 기술 스택
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "02", "기술 스택 & 인프라 구조")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 기술 스택 상세
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "기술 스택 상세", accent_color=C_PURPLE)

categories = [
    ("🖥️  Frontend", C_BLUE, [
        "Next.js 15 (App Router)",
        "TypeScript",
        "Tailwind CSS v4",
        "D3.js (지식 그래프)",
        "React Hooks (useState, useEffect, useRef)",
    ]),
    ("🗄️  Database", C_TEAL, [
        "Supabase (PostgreSQL)",
        "pgvector 확장 (벡터 검색)",
        "HNSW 인덱스 (ANN 검색)",
        "Row-Level Security",
        "REST API + RPC 함수",
    ]),
    ("🤖  AI / ML", C_GREEN, [
        "OpenAI text-embedding-3-small",
        "  → 1536차원, 임베딩 생성",
        "OpenAI gpt-4o-mini",
        "  → 답변·분류·요약·브리핑",
        "Ollama qwen2.5:7b (로컬 에이전트)",
    ]),
    ("⚙️  Backend", C_ORANGE, [
        "Next.js API Routes (서버리스)",
        "Python FastAPI (WebSocket 서버)",
        "NDJSON 스트리밍 응답",
        "Promise.allSettled 병렬 검색",
        "Vercel 엣지 캐시 (s-maxage=3600)",
    ]),
    ("🔐  Auth & Deploy", C_PURPLE, [
        "카카오 OAuth (Supabase provider)",
        "@supabase/ssr (서버 세션)",
        "Middleware (세션 갱신)",
        "Vercel (Next.js 배포)",
        "로컬 실행 (agents_server.py)",
    ]),
    ("📦  개발 도구", C_SLATE_400, [
        "generate_embeddings.py",
        "  → 임베딩 일괄 생성 스크립트",
        "build_knowledge_graph.py",
        "  → 지식 그래프 구축 스크립트",
        ".env.local (환경변수 관리)",
    ]),
]

for i, (title, col, items) in enumerate(categories):
    row = i // 3
    coli = i % 3
    lx = 0.25 + coli * 4.35
    ty = 1.15 + row * 3.05
    add_rounded_rect(slide, lx, ty, 4.15, 2.85, fill=C_SLATE_700, line=col)
    add_rect(slide, lx, ty, 4.15, 0.42, fill=col)
    add_text(slide, title, lx+0.1, ty+0.07, 3.9, 0.3, size=12, bold=True, color=C_WHITE)
    body = "\n".join(items)
    add_text(slide, body, lx+0.1, ty+0.5, 3.9, 2.2, size=9.5, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 섹션 헤더: 페이지 구조
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "03", "페이지 & 라우트 구조")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 페이지 라우트 맵
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "페이지 라우트 구조 (app/ 디렉토리)", accent_color=C_TEAL)

# 루트 노드
add_rounded_rect(slide, 5.7, 1.15, 2.0, 0.55, fill=C_BLUE, line=None)
add_text(slide, "app/ (루트)", 5.7, 1.15, 2.0, 0.55, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

routes = [
    # (label, lx, ty, w, color, desc)
    ("page.tsx\n🏠 홈", 0.2, 2.1, 1.65, C_BLUE, "DailyBriefing\n뉴스 목록 프리뷰"),
    ("news/\n📰 뉴스목록", 2.1, 2.1, 1.65, C_BLUE, "공개\n키워드 필터"),
    ("reports/\n📋 증권리포트", 4.0, 2.1, 1.65, C_PURPLE, "공개\nPDF 다운로드"),
    ("ask/\n💬 RAG Q&A", 5.9, 2.1, 1.65, C_GREEN, "RAG 챗봇\nAI 답변+출처"),
    ("graph/\n🕸️ 지식그래프", 7.8, 2.1, 1.65, C_TEAL, "D3.js 시각화\n엔티티 탐색"),
    ("insight/\n🔗 인사이트", 9.7, 2.1, 1.65, C_TEAL, "키워드/의미\n교차 참조"),
    ("agents/\n🤖 AI채팅", 11.6, 2.1, 1.65, C_ORANGE, "WebSocket\nlocalhost:8765"),
]

for label, lx, ty, w, col, desc in routes:
    add_rounded_rect(slide, lx, ty, w, 0.75, fill=C_SLATE_700, line=col)
    add_text(slide, label, lx, ty+0.04, w, 0.67, size=8.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx, ty+0.95, w, 0.5, size=7.5, color=C_SLATE_400, align=PP_ALIGN.CENTER)

# 하위 라우트
sub_routes = [
    ("report-analysis/\n📝 리포트 Pick", 0.2, 3.75, 2.5, C_PURPLE, "Notion 스타일\nHTML 리치 에디터"),
    ("telegram/\n📱 텔레그램", 2.9, 3.75, 2.2, C_TEAL, "채널별 메시지\n감성 분석"),
    ("board/\n💬 게시판", 5.3, 3.75, 2.0, C_SLATE_400, "CRUD\n이미지 업로드"),
    ("payment/\n💳 결제", 7.5, 3.75, 2.0, C_SLATE_600, "현재 미사용\n토스페이먼츠"),
    ("auth/callback\n🔐 OAuth 콜백", 9.7, 3.75, 2.2, C_ORANGE, "카카오 OAuth\n세션 생성"),
]
for label, lx, ty, w, col, desc in sub_routes:
    add_rounded_rect(slide, lx, ty, w, 0.75, fill=C_SLATE_700, line=col)
    add_text(slide, label, lx, ty+0.04, w, 0.67, size=8.5, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx, ty+0.95, w, 0.5, size=7.5, color=C_SLATE_400, align=PP_ALIGN.CENTER)

# API 라우트 섹션
add_text(slide, "API Routes (app/api/)", 0.2, 5.1, 4, 0.35, size=12, bold=True, color=C_ACCENT)
api_list = [
    ("chat/", "RAG Q&A", C_BLUE),
    ("briefing/", "일일 브리핑", C_ORANGE),
    ("graph/", "그래프 데이터+요약", C_PURPLE),
    ("insight/similar/", "의미 유사도", C_TEAL),
    ("board/", "게시판 CRUD", C_SLATE_400),
    ("report-pages/", "리포트 Pick", C_PURPLE),
    ("stocks/", "주가 데이터", C_GREEN),
    ("payment/confirm/", "결제 승인", C_SLATE_600),
    ("relative-performance/", "상대 수익률", C_GREEN),
]
ax, ay = 0.2, 5.5
for i, (route, desc, col) in enumerate(api_list):
    if i == 5:
        ax, ay = 6.8, 5.5
    add_rounded_rect(slide, ax, ay, 2.0, 0.45, fill=C_BG_DARKER, line=col)
    add_text(slide, f"/{route}", ax+0.08, ay+0.02, 1.2, 0.22, size=7.5, bold=True, color=col)
    add_text(slide, desc, ax+0.08, ay+0.23, 1.85, 0.2, size=7, color=C_SLATE_400)
    ax += 2.15
    if ax > 6.5 and ay == 5.5 and i < 5:
        ax = 0.2
        ay += 0.6


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 섹션 헤더: RAG Q&A
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "04", "RAG Q&A 챗봇 시스템")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RAG 파이프라인 다이어그램
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "RAG Q&A 챗봇 — 처리 파이프라인 (/ask + /api/chat)", accent_color=C_GREEN)

# STEP 박스들
steps = [
    ("①\n사용자 질문 입력", 0.25, 1.3, 1.9, C_SLATE_600),
    ("②\ngpt-4o-mini\n질문 분류", 2.45, 1.3, 2.0, C_ORANGE),
    ("③\ntext-embedding\n-3-small\n임베딩 생성", 4.75, 1.3, 2.0, C_BLUE),
    ("④\nSupabase\npgvector 검색", 7.05, 1.3, 2.0, C_TEAL),
    ("⑤\ngpt-4o-mini\n답변 생성", 9.35, 1.3, 2.0, C_GREEN),
    ("⑥\nNDJSON\n스트리밍\n출력", 11.65, 1.3, 1.45, C_PURPLE),
]

for label, lx, ty, w, col in steps:
    add_rounded_rect(slide, lx, ty, w, 1.15, fill=C_BG_DARKER, line=col)
    add_text(slide, label, lx, ty+0.1, w, 1.0, size=9.5, bold=True, color=col, align=PP_ALIGN.CENTER)

# 화살표
arrow_xs = [2.15, 4.45, 6.75, 9.05, 11.35]
for ax in arrow_xs:
    add_arrow(slide, ax, 1.87, ax+0.3, 1.87, color=C_SLATE_400, w=Pt(2))

# ② ③ 병렬 실행 표시
add_text(slide, "← Promise.all() 병렬 실행 →", 2.3, 2.55, 4.6, 0.3, size=8, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_rect(slide, 2.3, 2.42, 4.6, 0.06, fill=C_ORANGE)

# 분기 표시
add_rounded_rect(slide, 0.25, 2.65, 6.55, 1.8, fill=C_BG_DARKER, line=C_ORANGE)
add_text(slide, "🔀  검색 전략 분기 (isRecentQuery 결과)", 0.35, 2.72, 6.2, 0.3, size=10, bold=True, color=C_ORANGE)
add_rounded_rect(slide, 0.35, 3.08, 3.0, 1.25, fill=C_SLATE_700, line=C_BLUE)
add_text(slide, '📅  "recent" 분류\n최신 14일 검색', 0.35, 3.12, 3.0, 0.42, size=10, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "match_news_recent\nmatch_reports_recent\nmatch_telegrams_recent\nsince_days=14", 0.4, 3.55, 2.9, 0.72, size=8.5, color=C_SLATE_200)

add_rounded_rect(slide, 3.55, 3.08, 3.0, 1.25, fill=C_SLATE_700, line=C_TEAL)
add_text(slide, '🔍  "general" 분류\n전체 시맨틱 검색', 3.55, 3.12, 3.0, 0.42, size=10, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
add_text(slide, "match_news\nmatch_reports\nmatch_telegrams\n(HNSW 전체 탐색)", 3.6, 3.55, 2.9, 0.72, size=8.5, color=C_SLATE_200)

# 출력 형태
add_rounded_rect(slide, 6.8, 2.65, 6.25, 1.8, fill=C_BG_DARKER, line=C_GREEN)
add_text(slide, "📤  NDJSON 스트리밍 출력 형식", 6.95, 2.72, 5.9, 0.3, size=10, bold=True, color=C_GREEN)
add_text(slide, '{"type":"sources","news":[...],"reports":[...],"telegrams":[...],"isRecent":true}\n{"type":"text","data":"AI 답변 첫 번째 청크..."}\n{"type":"text","data":"...이어지는 텍스트..."}\n{"type":"error","message":"오류 메시지 (선택적)"}', 6.9, 3.1, 6.1, 1.3, size=8.5, color=C_SLATE_200)

# 핵심 설계 원칙
add_rounded_rect(slide, 0.25, 4.6, 12.8, 2.65, fill=C_SLATE_700, line=C_SLATE_600)
add_text(slide, "⚡  핵심 설계 원칙 & 트러블슈팅", 0.4, 4.68, 8, 0.32, size=12, bold=True, color=C_ACCENT)
principles = [
    ("임베딩 비대칭 문제", "짧은 질문과 긴 문서 임베딩의 코사인 유사도가 낮음 → threshold 필터 완전 제거", C_ORANGE),
    ("HNSW 인덱스 활용", "WHERE similarity > 0.3 추가 시 seq scan으로 fallback → 30k 행 풀스캔 타임아웃", C_RED),
    ("recent 모드 전략", "news_date_idx(B-tree) → 최근 200건 추출 → 그 안에서 유사도 정렬 (2단계)", C_TEAL),
    ("Promise.allSettled", "뉴스/리포트/텔레그램 검색 병렬 실행, 일부 실패해도 나머지 결과로 답변 가능", C_GREEN),
]
for j, (title, desc, col) in enumerate(principles):
    col_idx = j % 2
    row_idx = j // 2
    lx = 0.4 + col_idx * 6.4
    ty = 5.1 + row_idx * 0.88
    add_rect(slide, lx, ty, 0.05, 0.55, fill=col)
    add_text(slide, title, lx+0.15, ty, 2.0, 0.28, size=9.5, bold=True, color=col)
    add_text(slide, desc, lx+0.15, ty+0.28, 6.0, 0.28, size=8.5, color=C_SLATE_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RAG 검색 SQL 함수 구조
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "RAG — Supabase SQL 함수 구조", accent_color=C_GREEN)

# 일반 검색
add_rounded_rect(slide, 0.25, 1.1, 6.1, 3.1, fill=C_SLATE_700, line=C_TEAL)
add_text(slide, "🔍  일반 검색 함수 (전체 시맨틱)", 0.4, 1.18, 5.7, 0.35, size=12, bold=True, color=C_TEAL)
sql1 = """-- match_news / match_reports / match_telegrams
-- HNSW 인덱스 활용 (threshold 없음)

SELECT id, title, company, date::text,
       summary, keyword, link,
       1 - (embedding <=> query_embedding) AS similarity
FROM news
WHERE embedding IS NOT NULL
ORDER BY embedding <=> query_embedding
LIMIT match_count;

-- ✅ HNSW ivfflat → 빠른 ANN 검색
-- ✅ threshold 없음 → 임베딩 비대칭 문제 회피"""
add_text(slide, sql1, 0.35, 1.58, 5.9, 2.55, size=8.5, color=C_SLATE_200)

# 최근 14일 검색
add_rounded_rect(slide, 6.65, 1.1, 6.4, 3.1, fill=C_SLATE_700, line=C_ORANGE)
add_text(slide, "📅  최근 14일 검색 함수", 6.8, 1.18, 6.1, 0.35, size=12, bold=True, color=C_ORANGE)
sql2 = """-- match_news_recent (since_days=14)
-- 2단계 전략: 날짜 → 유사도

SELECT id, title, company, date::text,
       summary, keyword, link, similarity
FROM (
  SELECT n.*,
    1-(n.embedding <=> query_embedding) AS similarity
  FROM news n
  WHERE n.embedding IS NOT NULL
    AND n.date >= CURRENT_DATE
                 - (since_days||' days')::interval
  ORDER BY n.date DESC
  LIMIT 200          -- 날짜 인덱스로 빠르게 추출
) recent_sample
ORDER BY similarity DESC
LIMIT match_count;   -- 유사도 상위 N건 반환"""
add_text(slide, sql2, 6.75, 1.58, 6.2, 2.55, size=8.5, color=C_SLATE_200)

# 인덱스 설명
add_rounded_rect(slide, 0.25, 4.38, 6.1, 2.85, fill=C_BG_DARKER, line=C_BLUE)
add_text(slide, "🗂️  DB 인덱스 구조", 0.4, 4.48, 5.7, 0.35, size=12, bold=True, color=C_ACCENT)
idx_info = [
    ("HNSW 인덱스", "news.embedding, stock_reports.embedding, telegram_messages.embedding\nANN 검색 (Approximate Nearest Neighbor) 가속화", C_TEAL),
    ("news_date_idx", "CREATE INDEX ON news(date DESC)\nrecent 모드에서 30k행 풀스캔 타임아웃 방지 (B-tree)", C_ORANGE),
]
for j, (name, desc, col) in enumerate(idx_info):
    ty = 5.0 + j * 1.1
    add_rounded_rect(slide, 0.35, ty, 5.9, 0.95, fill=C_SLATE_700, line=col)
    add_text(slide, name, 0.5, ty+0.07, 2.5, 0.3, size=10, bold=True, color=col)
    add_text(slide, desc, 0.5, ty+0.4, 5.65, 0.48, size=8.5, color=C_SLATE_400)

# 분류기
add_rounded_rect(slide, 6.65, 4.38, 6.4, 2.85, fill=C_BG_DARKER, line=C_PURPLE)
add_text(slide, "🧠  질문 분류기 (isRecentQuery)", 6.8, 4.48, 6.1, 0.35, size=12, bold=True, color=RGBColor(0xc4, 0xb5, 0xfd))
classifier = """// gpt-4o-mini로 질문 분류 (max_tokens=5)
// 비용: ~$0.0001 per query

system: "최신 시황/동향/가격/실적 등
        '지금 현재' 정보가 중요하면 recent,
        기술 설명·역사·개념 등 시간에
        덜 민감하면 general로만 답해."

예시:
  "HBM 시장 현황은?" → recent (14일 검색)
  "HBM 기술 구조는?" → general (전체 검색)
  "삼성전자 최근 실적?" → recent
  "낸드 플래시 동작 원리?" → general"""
add_text(slide, classifier, 6.75, 4.88, 6.2, 2.28, size=8.5, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — 섹션 헤더: 일일 브리핑
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "05", "일일 브리핑 시스템")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 브리핑 시스템
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "일일 브리핑 — 데이터 흐름 & 캐싱 전략", accent_color=C_ORANGE)

# 데이터 수집 흐름
add_rounded_rect(slide, 0.25, 1.1, 12.8, 2.5, fill=C_SLATE_700, line=C_ORANGE)
add_text(slide, "📥  데이터 수집 (병렬, Promise.all)", 0.4, 1.18, 8, 0.32, size=12, bold=True, color=C_ORANGE)

sources_brief = [
    ("📰 뉴스 (10건)", "importance=eq.3\norder=date.desc\nlimit=10\n→ 중요도 최상위 최신 뉴스", C_BLUE, 0.4),
    ("📋 증권리포트 (5건)", "order=date.desc\nlimit=5\nsummary 필드 사용\n→ 최신 리포트 summary", C_PURPLE, 4.8),
    ("📱 텔레그램 (10건)", "order=date_utc.desc\n,forward_count.desc\nlimit=10\n→ 최신+바이럴 메시지", C_TEAL, 9.2),
]
for label, query, col, lx in sources_brief:
    add_rounded_rect(slide, lx, 1.58, 4.1, 1.88, fill=C_BG_DARK, line=col)
    add_text(slide, label, lx+0.1, 1.62, 3.9, 0.32, size=10, bold=True, color=col)
    add_text(slide, query, lx+0.1, 2.0, 3.9, 1.38, size=8.5, color=C_SLATE_200)

# 생성 흐름
add_rounded_rect(slide, 0.25, 3.75, 5.5, 2.25, fill=C_BG_DARKER, line=C_GREEN)
add_text(slide, "⚡  gpt-4o-mini 브리핑 생성", 0.4, 3.85, 5.2, 0.32, size=12, bold=True, color=C_GREEN)
gen_info = """출력 형식:
  📌 핵심 요약 (2~3문장)
  📈 주목 이슈 (3가지, 각 1~2문장)
  🔍 주목 키워드 (5~8개)

설정: max_tokens=500, temperature=0.3
비용: ~$0.003 per briefing"""
add_text(slide, gen_info, 0.4, 4.25, 5.2, 1.68, size=9.5, color=C_SLATE_200)

# 캐싱
add_rounded_rect(slide, 5.95, 3.75, 3.7, 2.25, fill=C_BG_DARKER, line=C_BLUE)
add_text(slide, "🚀  Vercel 엣지 캐시", 6.1, 3.85, 3.5, 0.32, size=12, bold=True, color=C_ACCENT)
cache_info = """Cache-Control:
  public, s-maxage=3600,
  stale-while-revalidate=7200

→ CDN에서 1시간 캐시
→ 대부분 요청: DB 접근 없음
→ 비용 최소화"""
add_text(slide, cache_info, 6.1, 4.25, 3.5, 1.68, size=9.5, color=C_SLATE_200)

# 수동 새로고침
add_rounded_rect(slide, 9.85, 3.75, 3.2, 2.25, fill=C_BG_DARKER, line=C_PURPLE)
add_text(slide, "🔄  수동 새로고침", 10.0, 3.85, 3.0, 0.32, size=12, bold=True, color=RGBColor(0xc4, 0xb5, 0xfd))
refresh_info = """캐시 우회 전략:

URL: /api/briefing
     ?t=${Date.now()}

fetch 옵션:
  { cache: "no-store" }

→ CDN 캐시 완전 무효화
→ 즉시 새 브리핑 생성"""
add_text(slide, refresh_info, 10.0, 4.25, 3.0, 1.68, size=9, color=C_SLATE_200)

# 컴포넌트
add_rounded_rect(slide, 0.25, 6.1, 12.8, 1.2, fill=C_SLATE_700, line=C_ORANGE)
add_text(slide, "🖥️  DailyBriefing 컴포넌트 (components/DailyBriefing.tsx)", 0.4, 6.18, 8, 0.32, size=11, bold=True, color=C_ORANGE)
add_text(slide, "상태: loading / briefing / date / error / refreshing  ·  초기 로드 후 1시간 캐시 응답  ·  새로고침 버튼으로 즉시 재생성 가능", 0.4, 6.57, 12.4, 0.6, size=9.5, color=C_SLATE_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 섹션 헤더: AI 에이전트
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "06", "AI 에이전트 시황 채팅")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — AI 에이전트 구조
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "AI 에이전트 시황 채팅 — agents_server.py 구조", accent_color=C_GREEN)

# 서버 정보
add_rounded_rect(slide, 0.25, 1.1, 4.0, 1.5, fill=C_SLATE_700, line=C_GREEN)
add_text(slide, "🐍  FastAPI WebSocket 서버", 0.4, 1.18, 3.7, 0.32, size=11, bold=True, color=C_GREEN)
add_text(slide, "포트: 8765\n모델: Ollama qwen2.5:7b\n실행: python agents_server.py", 0.4, 1.55, 3.7, 0.98, size=9.5, color=C_SLATE_200)

add_rounded_rect(slide, 4.45, 1.1, 4.2, 1.5, fill=C_SLATE_700, line=C_BLUE)
add_text(slide, "🌐  연결 구조", 4.6, 1.18, 3.9, 0.32, size=11, bold=True, color=C_ACCENT)
add_text(slide, "브라우저 → ws://localhost:8765/ws\nVercel 배포 불가 → 로컬 직접 실행\n시작 스크립트: .\\start_agents.ps1", 4.6, 1.55, 3.9, 0.98, size=9.5, color=C_SLATE_200)

add_rounded_rect(slide, 8.85, 1.1, 4.2, 1.5, fill=C_SLATE_700, line=C_ORANGE)
add_text(slide, "📊  데이터 로드", 9.0, 1.18, 3.9, 0.32, size=11, bold=True, color=C_ORANGE)
add_text(slide, "시작 시: 뉴스 30, 텔레그램 30\n         리포트 20, 분석 12건\n30틱(~5분)마다 자동 갱신", 9.0, 1.55, 3.9, 0.98, size=9.5, color=C_SLATE_200)

# 에이전트 5개
add_text(slide, "🤖  에이전트 5개 — 각기 다른 관점으로 토론", 0.25, 2.78, 8, 0.35, size=12, bold=True, color=C_ACCENT)
agents = [
    ("🐂 강세론자\nbull", "news, reports", "낙관적 시장 뷰\n상승 근거 중심", C_GREEN),
    ("🐻 약세론자\nbear", "telegram, analysis", "비관적 시장 뷰\n리스크 강조", C_RED),
    ("⚠️ 리스크\nrisk", "telegram, news", "위험 요인 탐색\n불확실성 분석", C_ORANGE),
    ("📊 애널리스트\nanalyst", "reports, analysis", "데이터 기반 분석\n객관적 수치 인용", C_BLUE),
    ("🌍 매크로\nmacro", "news, analysis", "거시경제 관점\n글로벌 이슈 연결", C_PURPLE),
]
for j, (name, focus, desc, col) in enumerate(agents):
    lx = 0.25 + j * 2.58
    add_rounded_rect(slide, lx, 3.2, 2.42, 2.0, fill=C_BG_DARKER, line=col)
    add_rect(slide, lx, 3.2, 2.42, 0.4, fill=col)
    add_text(slide, name, lx, 3.2, 2.42, 0.4, size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, f"소스: {focus}", lx+0.1, 3.65, 2.2, 0.3, size=8, color=C_SLATE_400)
    add_text(slide, desc, lx+0.1, 3.98, 2.2, 0.55, size=8.5, color=C_SLATE_200)

# 동작 로직
add_rounded_rect(slide, 0.25, 5.35, 12.8, 1.95, fill=C_SLATE_700, line=C_SLATE_600)
add_text(slide, "⚙️  대화 구동 로직", 0.4, 5.43, 5, 0.32, size=12, bold=True, color=C_ACCENT)
logic_items = [
    ("순환 발언", "에이전트 5개가 순서대로 Ollama로 메시지 생성 → WebSocket broadcast"),
    ("대화 환기", "15~20틱마다 '현재 시황이 어떻지?' 자동 주입 → 대화 방향 유지"),
    ("버스트 모드", "논쟁적 발언 감지 시 1명 추가 반응 → 더 활발한 토론"),
    ("한국어 보장", "CJK 문자 감지 시 1회 재시도 후 건너뜀 (has_cjk() 정규식)"),
    ("타임아웃", "Ollama 요청 35초, num_predict: 200으로 응답 길이 제한"),
]
ly = 5.82
for j, (key, val) in enumerate(logic_items):
    col_idx = j % 3
    row_idx = j // 3
    lx = 0.38 + col_idx * 4.2
    ty = ly + row_idx * 0.52
    add_rect(slide, lx, ty+0.06, 0.05, 0.3, fill=C_GREEN)
    add_text(slide, f"{key}: ", lx+0.12, ty, 1.0, 0.32, size=8.5, bold=True, color=C_GREEN)
    add_text(slide, val, lx+0.12, ty+0.25, 3.9, 0.28, size=8, color=C_SLATE_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — 섹션 헤더: 지식 그래프
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "07", "지식 그래프 & 인사이트 시스템")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — 지식 그래프 & 인사이트
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "지식 그래프 & 인사이트 — 엔티티 관계 탐색", accent_color=C_PURPLE)

# 지식 그래프
add_rounded_rect(slide, 0.25, 1.1, 7.9, 6.15, fill=C_SLATE_700, line=C_PURPLE)
add_text(slide, "🕸️  지식 그래프 (/graph) — D3.js Force-Directed Graph", 0.4, 1.18, 7.6, 0.35, size=12, bold=True, color=RGBColor(0xc4, 0xb5, 0xfd))

# 노드 타입
node_types = [
    ("기업", "삼성전자, SK하이닉스, 엔비디아...", C_BLUE),
    ("제품/기술", "HBM, NAND, CXL, LPDDR...", C_GREEN),
    ("지표", "ASP, 영업이익, 점유율...", C_ORANGE),
    ("이벤트", "실적발표, 공급계약, 규제...", C_RED),
    ("섹터", "메모리, 파운드리, 장비...", C_PURPLE),
]
add_text(slide, "노드 타입 (엔티티 2,121개):", 0.4, 1.6, 5, 0.28, size=10, bold=True, color=C_SLATE_400)
for j, (ntype, examples, col) in enumerate(node_types):
    ty = 1.95 + j * 0.52
    add_rounded_rect(slide, 0.4, ty, 0.9, 0.38, fill=col)
    add_text(slide, ntype, 0.4, ty, 0.9, 0.38, size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, examples, 1.4, ty+0.04, 6.4, 0.3, size=8.5, color=C_SLATE_400)

add_text(slide, "인터랙션:", 0.4, 4.6, 3, 0.28, size=10, bold=True, color=C_SLATE_400)
interactions = [
    "줌 / 드래그 / 노드 클릭",
    "타입 필터 & 최소 관계 강도 슬라이더",
    "엔티티 검색 (이름으로 필터)",
    "클릭 시 최신 뉴스·리포트·텔레그램 표시",
    "AI 시황 요약 스트리밍 (graph/summary API)",
]
for j, item in enumerate(interactions):
    add_text(slide, f"• {item}", 0.4, 4.93 + j*0.44, 7.5, 0.38, size=9.5, color=C_SLATE_200)

add_text(slide, "노드 크기 = 멘션 수  ·  엣지 두께 = 관계 강도 (공동 출현 횟수)", 0.4, 7.05, 7.6, 0.12, size=8.5, color=C_SLATE_600)

# 인사이트
add_rounded_rect(slide, 8.4, 1.1, 4.65, 6.15, fill=C_SLATE_700, line=C_TEAL)
add_text(slide, "🔗  인사이트 (/insight)", 8.55, 1.18, 4.4, 0.35, size=12, bold=True, color=C_TEAL)

modes = [
    ("키워드 매칭", "공통 키워드 수로 관련도 계산\nnews(keyword), reports(keyword),\ntelegram(keywords)\n→ 정확하지만 표기 불일치 시 누락", C_ORANGE),
    ("의미 유사도", "pgvector 코사인 유사도\nthreshold=0.1 (비대칭 임베딩 고려)\n교차 타입 검색 가능\n→ 표기 달라도 의미 연결", C_TEAL),
]
for j, (mode, desc, col) in enumerate(modes):
    ty = 1.65 + j * 2.15
    add_rounded_rect(slide, 8.55, ty, 4.35, 1.95, fill=C_BG_DARKER, line=col)
    add_text(slide, mode, 8.7, ty+0.1, 4.0, 0.32, size=11, bold=True, color=col)
    add_text(slide, desc, 8.7, ty+0.48, 4.0, 1.38, size=9, color=C_SLATE_200)

add_text(slide, "데이터 구조:", 8.55, 5.97, 3, 0.28, size=10, bold=True, color=C_SLATE_400)
db_info = [
    "entities: 2,121개 (노드)",
    "entity_mentions: 47,609건 (문서↔엔티티)",
    "entity_relations: 9,537건 (엔티티↔엔티티)",
    "구축: build_knowledge_graph.py",
    "  → gpt-4o-mini, 최근 6개월 데이터",
]
add_text(slide, "\n".join(db_info), 8.55, 6.3, 4.35, 0.88, size=8.5, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — 섹션 헤더: DB 구조
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "08", "Supabase DB 구조 & SQL 함수")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — DB 테이블 구조
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "Supabase DB 테이블 구조", accent_color=C_TEAL)

tables = [
    ("news", "뉴스 기사", C_BLUE, [
        "id (int, PK)",
        "title (text)",
        "company (text)",
        "date (date)",
        "summary (text)",
        "keyword (text)",
        "link (text)",
        "importance (int) — 1~3",
        "embedding (vector 1536)",
    ]),
    ("stock_reports", "증권사 리포트", C_PURPLE, [
        "id (int, PK)",
        "title (text)",
        "securities_firm (text)",
        "date (date)",
        "summary (text)  ← 주의",
        "keyword (text)",
        "link (text)",
        "embedding (vector 1536)",
    ]),
    ("telegram_messages", "텔레그램 메시지", C_TEAL, [
        "id (int, PK)",
        "channel (text)",
        "summary (text)",
        "date_utc (timestamptz)",
        "sentiment (text)",
        "keywords (text)",
        "forward_count (int)",
        "embedding (vector 1536)",
    ]),
    ("entities", "지식 그래프 노드", C_GREEN, [
        "id (int, PK)",
        "name (text)",
        "type (text)",
        "  company/product/metric/",
        "  event/sector",
        "mention_count (int)",
        "description (text)",
    ]),
    ("entity_mentions", "문서↔엔티티 연결", C_ORANGE, [
        "id (int, PK)",
        "entity_id (int, FK→entities)",
        "source_type (text)",
        "  news/report/telegram",
        "source_id (int)",
        "date (date)",
    ]),
    ("entity_relations", "엔티티↔엔티티 관계", C_RED, [
        "id (int, PK)",
        "entity1_id (int, FK→entities)",
        "entity2_id (int, FK→entities)",
        "weight (int) — 공동출현 횟수",
        "relation_type (text)",
    ]),
]

for i, (name, desc, col, fields) in enumerate(tables):
    row = i // 3
    col_idx = i % 3
    lx = 0.25 + col_idx * 4.35
    ty = 1.1 + row * 3.2
    add_rounded_rect(slide, lx, ty, 4.15, 3.05, fill=C_SLATE_700, line=col)
    add_rect(slide, lx, ty, 4.15, 0.42, fill=col)
    add_text(slide, name, lx+0.1, ty+0.07, 2.5, 0.28, size=11, bold=True, color=C_WHITE)
    add_text(slide, desc, lx+2.5, ty+0.1, 1.5, 0.25, size=8.5, color=C_WHITE, align=PP_ALIGN.RIGHT)
    field_text = "\n".join(fields)
    add_text(slide, field_text, lx+0.1, ty+0.5, 3.9, 2.48, size=8.5, color=C_SLATE_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — SQL 함수 & 인덱스
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "Supabase SQL 함수 목록 & 인덱스", accent_color=C_TEAL)

funcs = [
    ("match_news", "전체 뉴스 시맨틱 검색", "HNSW ORDER BY embedding <=> query LIMIT n", C_BLUE),
    ("match_reports", "전체 리포트 시맨틱 검색", "HNSW ORDER BY embedding <=> query LIMIT n", C_PURPLE),
    ("match_telegrams", "전체 텔레그램 시맨틱 검색", "HNSW ORDER BY embedding <=> query LIMIT n", C_TEAL),
    ("match_news_recent", "최근 14일 뉴스 검색", "날짜 인덱스로 200건 → 유사도 정렬 LIMIT n", C_BLUE),
    ("match_reports_recent", "최근 14일 리포트 검색", "날짜 필터 200건 → 유사도 정렬 LIMIT n", C_PURPLE),
    ("match_telegrams_recent", "최근 14일 텔레그램 검색", "날짜 필터 200건 → 유사도 정렬 LIMIT n", C_TEAL),
    ("get_entity_docs", "엔티티별 최신 문서 조회", "날짜 DESC 정렬, 뉴스/리포트/텔레그램 각 최신", C_GREEN),
]

add_text(slide, "🔧  RPC 함수 목록", 0.25, 1.1, 5, 0.32, size=13, bold=True, color=C_ACCENT)
for j, (name, desc, impl, col) in enumerate(funcs):
    ty = 1.5 + j * 0.76
    add_rounded_rect(slide, 0.25, ty, 12.8, 0.68, fill=C_BG_DARKER, line=col)
    add_rect(slide, 0.25, ty, 0.06, 0.68, fill=col)
    add_text(slide, name, 0.4, ty+0.05, 2.8, 0.28, size=10, bold=True, color=col)
    add_text(slide, desc, 0.4, ty+0.38, 2.8, 0.25, size=8.5, color=C_SLATE_400)
    add_text(slide, impl, 3.4, ty+0.18, 9.4, 0.32, size=9, color=C_SLATE_200)

# 인덱스
add_text(slide, "🗂️  DB 인덱스", 0.25, 6.9, 5, 0.3, size=13, bold=True, color=C_ACCENT)
indexes = [
    ("news_date_idx", "B-tree", "news(date DESC)", "recent 모드 30k행 풀스캔 방지, 최근 200건 추출에 사용"),
    ("HNSW 인덱스 ×3", "HNSW (pgvector)", "*.embedding", "ANN 검색 가속, WHERE threshold 필터 추가 시 비활성화 주의"),
]
ix = 0.25
for name, itype, col_t, desc in indexes:
    add_rounded_rect(slide, ix, 7.25, 6.25, 0.8, fill=C_SLATE_700, line=C_SLATE_600)
    add_text(slide, name, ix+0.1, 7.3, 2.5, 0.28, size=10, bold=True, color=C_ACCENT)
    add_text(slide, f"타입: {itype}  |  대상: {col_t}", ix+0.1, 7.6, 5.9, 0.25, size=8.5, color=C_SLATE_400)
    add_text(slide, desc, ix+0.1, 7.88, 5.9, 0.12, size=7.5, color=C_SLATE_600)
    ix += 6.5


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — 섹션 헤더: 데이터 현황
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
section_header(slide, "09", "데이터 현황 & 비용 구조")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — 데이터 현황 & 비용
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
slide_header(slide, "데이터 현황 & 운영 비용 구조", accent_color=C_GREEN)

# 데이터 현황 테이블
add_rounded_rect(slide, 0.25, 1.1, 7.8, 4.2, fill=C_SLATE_700, line=C_TEAL)
add_text(slide, "📊  데이터 현황 (2026-05-25 기준)", 0.4, 1.18, 7.5, 0.32, size=12, bold=True, color=C_TEAL)

headers = ["소스", "테이블", "건수", "임베딩"]
col_xs = [0.4, 1.8, 4.2, 5.7]
add_rect(slide, 0.25, 1.55, 7.8, 0.38, fill=C_SLATE_600)
for h, lx in zip(headers, col_xs):
    add_text(slide, h, lx, 1.6, 1.3, 0.28, size=9, bold=True, color=C_WHITE)

data_rows = [
    ("뉴스 기사", "news", "~30,000건", "✅ 1536dim", C_BLUE),
    ("증권 리포트", "stock_reports", "~584건", "✅ 1536dim", C_PURPLE),
    ("텔레그램", "telegram_messages", "~5,558건", "✅ 1536dim", C_TEAL),
    ("지식 그래프 엔티티", "entities", "2,121개", "—", C_GREEN),
    ("엔티티 멘션", "entity_mentions", "47,609건", "—", C_ORANGE),
    ("엔티티 관계", "entity_relations", "9,537건", "—", C_RED),
    ("리포트 Pick 페이지", "report_pages", "운영 중", "—", C_PURPLE),
    ("게시판", "board_posts", "운영 중", "—", C_SLATE_400),
]
for j, (src, table, count, emb, col) in enumerate(data_rows):
    ty = 1.98 + j * 0.4
    bg = C_BG_DARK if j % 2 == 0 else C_BG_DARKER
    add_rect(slide, 0.25, ty, 7.8, 0.38, fill=bg)
    add_rect(slide, 0.25, ty, 0.05, 0.38, fill=col)
    add_text(slide, src, 0.4, ty+0.04, 1.35, 0.3, size=8.5, color=C_SLATE_200)
    add_text(slide, table, 1.8, ty+0.04, 2.35, 0.3, size=8, color=C_SLATE_400)
    add_text(slide, count, 4.2, ty+0.04, 1.4, 0.3, size=9, bold=True, color=col)
    add_text(slide, emb, 5.7, ty+0.04, 2.2, 0.3, size=8.5, color=C_GREEN if "✅" in emb else C_SLATE_600)

# 비용
add_rounded_rect(slide, 8.3, 1.1, 4.75, 4.2, fill=C_SLATE_700, line=C_GREEN)
add_text(slide, "💰  운영 비용 구조", 8.45, 1.18, 4.4, 0.32, size=12, bold=True, color=C_GREEN)
costs = [
    ("Vercel", "무료 플랜 (Hobby)", C_BLUE, "서버리스 함수 + CDN"),
    ("Supabase", "Free tier\n→ 500MB DB, 1GB 파일", C_TEAL, "PostgreSQL + pgvector"),
    ("OpenAI\n임베딩", "~$0.0001 / 1K tokens\n초기 구축 시 일회성", C_GREEN, "text-embedding-3-small"),
    ("OpenAI\n답변", "~$0.003 / Q&A 1회\n~$0.003 / 브리핑 1회", C_ORANGE, "gpt-4o-mini"),
    ("Ollama", "로컬 실행 → 무료\n(GPU 없이도 동작)", C_PURPLE, "qwen2.5:7b"),
]
for j, (name, cost, col, desc) in enumerate(costs):
    ty = 1.6 + j * 0.7
    add_rounded_rect(slide, 8.45, ty, 4.45, 0.6, fill=C_BG_DARKER, line=col)
    add_text(slide, name, 8.55, ty+0.04, 1.2, 0.52, size=9, bold=True, color=col)
    add_text(slide, cost, 9.85, ty+0.04, 2.0, 0.52, size=8.5, color=C_SLATE_200)
    add_text(slide, desc, 11.9, ty+0.14, 0.9, 0.3, size=7, color=C_SLATE_600, align=PP_ALIGN.RIGHT)

# 환경변수
add_rounded_rect(slide, 0.25, 5.45, 12.8, 1.85, fill=C_BG_DARKER, line=C_SLATE_600)
add_text(slide, "🔑  환경변수 (.env.local)", 0.4, 5.53, 5, 0.32, size=12, bold=True, color=C_ACCENT)
envs = [
    ("NEXT_PUBLIC_SUPABASE_URL", "Supabase 프로젝트 URL"),
    ("NEXT_PUBLIC_SUPABASE_ANON_KEY", "Supabase 공개 키"),
    ("OPENAI_API_KEY", "GPT-4o-mini + 임베딩"),
    ("NEXT_PUBLIC_KAKAO_CLIENT_ID", "카카오 OAuth"),
    ("NEXT_PUBLIC_ADMIN_EMAIL", "어드민 이메일"),
    ("NEXT_PUBLIC_TOSS_CLIENT_KEY", "토스페이먼츠 (미사용)"),
]
ex, ey = 0.4, 5.92
for j, (key, desc) in enumerate(envs):
    col_idx = j % 3
    row_idx = j // 3
    lx = ex + col_idx * 4.22
    ty = ey + row_idx * 0.45
    add_text(slide, f"{key}", lx, ty, 2.6, 0.28, size=8, bold=True, color=C_ACCENT)
    add_text(slide, f"  {desc}", lx, ty+0.22, 2.6, 0.22, size=7.5, color=C_SLATE_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — 마무리 / 요약
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_BG_DARKER)
add_rect(slide, 0, 0, 13.33, 0.08, fill=C_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill=C_BLUE)

add_text(slide, "시스템 요약", 0.5, 0.7, 12.33, 0.55, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

summary_items = [
    ("🗄️  데이터 통합", "뉴스 30K + 리포트 584 + 텔레그램 5.5K\npgvector 1536dim 임베딩 저장", C_TEAL),
    ("🤖  AI 활용", "gpt-4o-mini: Q&A · 브리핑 · 엔티티 분류 · 요약\nOllama qwen2.5:7b: 실시간 에이전트 토론", C_GREEN),
    ("⚡  성능 최적화", "HNSW 인덱스 (ANN) + news_date_idx (B-tree)\nrecent 모드 타임아웃 없이 14일 검색", C_ORANGE),
    ("🌐  배포 구조", "Next.js → Vercel (서버리스)\nagents_server.py → 사용자 로컬 실행", C_BLUE),
    ("📊  핵심 기능", "RAG Q&A · 일일 브리핑 · 지식 그래프\n인사이트 교차 참조 · AI 에이전트 채팅", C_PURPLE),
    ("💡  설계 원칙", "임베딩 비대칭 → threshold 제거\nPromise.allSettled → 부분 실패 허용", C_RED),
]

for i, (title, desc, col) in enumerate(summary_items):
    col_idx = i % 3
    row_idx = i // 3
    lx = 0.4 + col_idx * 4.3
    ty = 1.5 + row_idx * 2.5
    add_rounded_rect(slide, lx, ty, 4.1, 2.2, fill=C_SLATE_700, line=col)
    add_rect(slide, lx, ty, 4.1, 0.08, fill=col)
    add_text(slide, title, lx+0.15, ty+0.18, 3.8, 0.4, size=13, bold=True, color=col)
    add_text(slide, desc, lx+0.15, ty+0.68, 3.8, 1.38, size=10, color=C_SLATE_200)

add_text(slide, "next.js 15  ·  Supabase pgvector  ·  OpenAI  ·  D3.js  ·  FastAPI  ·  Ollama",
         0.5, 7.1, 12.33, 0.3, size=9, color=C_SLATE_600, align=PP_ALIGN.CENTER)


# 저장
out = r"c:\Users\yh900\SemiCon\SemiCon_Architecture.pptx"
prs.save(out)
print(f"PPT saved: {out}")
print(f"Slides: {len(prs.slides)}")
