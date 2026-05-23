from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

prs = Presentation()
prs.slide_width  = Cm(33.87)   # 16:9 widescreen
prs.slide_height = Cm(19.05)

W = prs.slide_width
H = prs.slide_height

# ── colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x23, 0x44)
GOLD    = RGBColor(0xF0, 0xA5, 0x00)
GOLD_LT = RGBColor(0xFF, 0xD0, 0x60)
OFFWHITE= RGBColor(0xF7, 0xF8, 0xFC)
SLATE   = RGBColor(0x2C, 0x3E, 0x5A)
GRAY    = RGBColor(0x6B, 0x7B, 0x93)
ROW_ODD = RGBColor(0xEF, 0xF2, 0xF8)
ROW_EVN = RGBColor(0xFA, 0xFB, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

def rgb(color):
    from pptx.oxml.ns import qn
    from lxml import etree
    return color

def bg(slide, color):
    from pptx.oxml.ns import qn
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def tb(slide, text, x, y, w, h, size, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    run.font.name  = "Pretendard" if False else "맑은 고딕"
    return txb

def header_bar(slide, title, subtitle=""):
    # navy header band
    rect(slide, 0, 0, W, Cm(2.8), NAVY)
    # gold accent line at bottom of header
    rect(slide, 0, Cm(2.8), W, Cm(0.12), GOLD)
    tb(slide, title,
       Cm(1.2), Cm(0.3), W - Cm(2.4), Cm(2.0),
       22, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle,
           Cm(1.2), Cm(2.0), W - Cm(2.4), Cm(0.7),
           13, color=GOLD_LT, italic=True)

def footer_note(slide, note):
    # gold left accent
    rect(slide, Cm(1.0), H - Cm(1.3), Cm(0.25), Cm(0.85), GOLD)
    tb(slide, note,
       Cm(1.5), H - Cm(1.35), W - Cm(3), Cm(0.9),
       11, italic=True, color=SLATE)

# ── slide builders ──────────────────────────────────────────────────────────

def title_slide(title, subtitle=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg(sl, NAVY)
    # right geometric panel
    rect(sl, W * 0.62, 0, W * 0.38, H, SLATE)
    rect(sl, W * 0.62, 0, Cm(0.35), H, GOLD)
    # gold bars
    rect(sl, 0, Cm(1.5), W * 0.62, Cm(0.2), GOLD)
    rect(sl, 0, H - Cm(1.8), W * 0.62, Cm(0.2), GOLD_LT)
    # title
    tb(sl, title,
       Cm(1.5), Cm(4.5), W * 0.58, Cm(6),
       44, bold=True, color=WHITE)
    if subtitle:
        tb(sl, subtitle,
           Cm(1.5), Cm(11.5), W * 0.58, Cm(2.5),
           20, color=GOLD_LT, italic=True)
    # decorative dot grid (right panel label)
    tb(sl, "MEMORY INDUSTRY\nANALYSIS",
       W * 0.64, Cm(8), W * 0.34, Cm(3),
       16, bold=True, color=GOLD_LT, align=PP_ALIGN.CENTER)
    return sl

def section_slide(num, title):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, NAVY)
    # vertical gold bar left
    rect(sl, Cm(1.0), Cm(2.5), Cm(0.5), Cm(14), GOLD)
    # chapter number
    tb(sl, f"Chapter {num:02d}",
       Cm(2.2), Cm(4), Cm(20), Cm(1.5),
       18, color=GOLD)
    # title
    tb(sl, title,
       Cm(2.2), Cm(5.5), W - Cm(4), Cm(5),
       40, bold=True, color=WHITE)
    # gold underline
    rect(sl, Cm(2.2), Cm(11.5), Cm(12), Cm(0.18), GOLD)
    return sl

def content_slide(title, bullets, note="", subtitle=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, OFFWHITE)
    header_bar(sl, title, subtitle)
    if note:
        footer_note(sl, note)

    tf_h = H - Cm(2.8) - (Cm(1.6) if note else Cm(0.4))
    txb = slide_textbox(sl, Cm(1.2), Cm(3.1), W - Cm(2.4), tf_h)

    for item in bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        add_bullet(txb.text_frame, text, level)
    return sl

def slide_textbox(sl, x, y, w, h):
    txb = sl.shapes.add_textbox(x, y, w, h)
    txb.text_frame.word_wrap = True
    # clear default empty paragraph
    return txb

def add_bullet(tf, text, level=0):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    p = tf.add_paragraph()
    p.level = level

    run = p.add_run()
    run.text = text

    if level == 0:
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = NAVY
        p.space_before = Pt(8)
        p.space_after  = Pt(2)
        # bullet symbol
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', '맑은 고딕')
    elif level == 1:
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.color.rgb = SLATE
        p.space_before = Pt(3)
        p.space_after  = Pt(1)
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '–')
    else:
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.italic = True
        run.font.color.rgb = GRAY
        p.space_before = Pt(2)
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '·')

    run.font.name = "맑은 고딕"


def table_slide(title, headers, rows, note="", subtitle=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, OFFWHITE)
    header_bar(sl, title, subtitle)
    if note:
        footer_note(sl, note)

    t_top  = Cm(3.2)
    t_bot  = H - (Cm(1.8) if note else Cm(0.5))
    t_h    = t_bot - t_top
    cols   = len(headers)
    col_w  = (W - Cm(2.4)) / cols

    tbl = sl.shapes.add_table(
        len(rows) + 1, cols,
        Cm(1.2), t_top, W - Cm(2.4), t_h
    ).table

    # set column widths
    for ci in range(cols):
        tbl.columns[ci].width = int(col_w)

    # header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold  = True
        run.font.size  = Pt(13)
        run.font.color.rgb = WHITE
        run.font.name  = "맑은 고딕"

    # data rows
    for ri, row in enumerate(rows):
        fill_color = ROW_ODD if ri % 2 == 0 else ROW_EVN
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size  = Pt(12)
            run.font.color.rgb = NAVY if ci == 0 else SLATE
            run.font.bold  = (ci == 0)
            run.font.name  = "맑은 고딕"

    return sl


def two_col_slide(title, left_bullets, right_bullets, note="", subtitle=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, OFFWHITE)
    header_bar(sl, title, subtitle)
    if note:
        footer_note(sl, note)

    col_w = (W - Cm(3.6)) / 2
    col_h = H - Cm(2.8) - (Cm(1.6) if note else Cm(0.4))

    # left column
    txb_l = slide_textbox(sl, Cm(1.2), Cm(3.1), col_w, col_h)
    for item in left_bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        add_bullet(txb_l.text_frame, text, level)

    # divider line
    rect(sl, Cm(1.2) + col_w + Cm(0.5), Cm(3.1), Cm(0.06), col_h, GOLD)

    # right column
    txb_r = slide_textbox(sl, Cm(1.2) + col_w + Cm(1.2), Cm(3.1), col_w, col_h)
    for item in right_bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        add_bullet(txb_r.text_frame, text, level)

    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════

# ── Title ──────────────────────────────────────────────────────────────────
title_slide("메모리 산업 분석", "반도체 뉴스레터 리포트 · 2025")

# ═══════════════════════════════════════════════════════════════════════════
# Ch1. 메모리 산업의 정의와 원리
# ═══════════════════════════════════════════════════════════════════════════
section_slide(1, "메모리 산업의\n정의와 원리")

content_slide(
    "메모리 반도체란?",
    [
        (0, "데이터를 저장하는 반도체 — DRAM과 NAND가 핵심"),
        (1, "DRAM: 빠르지만 전원 끄면 사라지는 휘발성 메모리"),
        (1, "NAND Flash: 느리지만 전원 없이도 데이터 유지 (비휘발성)"),
        (0, "메모리는 범용(Commodity) 제품"),
        (1, "동일 규격이면 어느 제조사 제품이든 호환 → 가격 경쟁이 치열"),
        (1, "수요·공급 불균형이 곧 가격 사이클 형성"),
        (0, "로직 반도체(CPU·GPU)와의 차이"),
        (1, "로직은 설계 차별화로 마진 보호 가능"),
        (1, "메모리는 공정 미세화(원가 절감)가 유일한 경쟁 수단"),
    ],
    note="메모리 = 범용 부품 → 가격 사이클의 출발점"
)

table_slide(
    "DRAM vs NAND 비교",
    ["구분", "DRAM", "NAND Flash"],
    [
        ["역할",       "주기억장치 (RAM)",          "보조기억장치 (SSD·eMMC)"],
        ["휘발성",     "휘발성 (전원 차단 시 소멸)", "비휘발성 (전원 무관)"],
        ["속도",       "수 ns (매우 빠름)",          "수십 μs (상대적으로 느림)"],
        ["집적도 향상","셀 크기 축소 (2D)",          "3D 적층 (레이어 수 증가)"],
        ["주요 수요처","서버·PC·모바일",             "SSD·스마트폰·데이터센터"],
        ["가격 변동성","높음",                        "더 높음"],
    ],
    note="DRAM과 NAND는 수요처·기술 방향이 다르지만 같은 사이클 압력에 노출"
)

content_slide(
    "Commodity 특성과 가격 변동성",
    [
        (0, "메모리 = 대표적 Commodity 반도체"),
        (1, "가격 탄력성 높음 — 가격 하락 시 수요 급증, 상승 시 수요 위축"),
        (1, "재고 증가 → 가격 급락 → 감산 → 재고 소진 → 가격 급등 반복"),
        (0, "공급자 교섭력이 수요자보다 낮은 시장"),
        (1, "고객사는 복수 소싱 → 벤더 경쟁 심화"),
        (1, "장기계약(LTA)으로 일부 완충하나 스팟 가격 영향 불가피"),
        (0, "투자 관점"),
        (1, "가격 사이클 파악이 투자 수익의 핵심 드라이버"),
        (1, "재고 수준 + Bit Growth 지표가 선행지표"),
    ],
    note="Commodity 속성 → 가격 변동성 → 사이클 투자 기회"
)

content_slide(
    "Wafer ≠ Bit — 원가 구조의 핵심",
    [
        (0, "웨이퍼 투입량(Wafer)과 생산 비트(Bit)는 다른 개념"),
        (1, "같은 웨이퍼에서 공정 미세화·레이어 증가로 Bit 생산량 증가"),
        (1, "웨이퍼 Capa를 늘리지 않아도 Bit 공급 확대 가능"),
        (0, "Bit Growth = 기술 향상에 의한 자연 공급 증가"),
        (1, "DRAM: 연간 약 15~20% Bit Growth 전통적으로 유지"),
        (1, "NAND: 3D 레이어 증가로 더 큰 Bit Growth (20~30%)"),
        (0, "투자 관점"),
        (1, "웨이퍼 Capa 동결해도 Bit 공급이 늘 수 있어 공급 과잉 주의"),
        (1, "Bit Growth 둔화 = 공급 부족 신호"),
    ],
    note="Wafer Capa ≠ Bit 공급 — 기술 미세화가 핵심 변수"
)

content_slide(
    "Re-entrant Flow — 메모리 제조의 특수성",
    [
        (0, "메모리 공정은 Re-entrant Flow 구조"),
        (1, "동일 장비에 웨이퍼가 수십 회 반복 투입 (선형 공정이 아님)"),
        (1, "장비 가동률 변화가 생산량에 비선형적으로 영향"),
        (0, "가동률 하락의 불균형 효과"),
        (1, "가동률 10% 하락 → 생산량 10% 이상 감소 가능"),
        (1, "병목 장비 한 대가 전체 라인 속도를 결정"),
        (0, "투자 관점"),
        (1, "감산 발표 후 실제 Bit 감소까지 수개월 시차 발생"),
        (1, "감산 규모 × Re-entrant 비선형성 = 실제 공급 충격 과소평가 금물"),
    ],
    note="Re-entrant Flow → 감산 효과가 예상보다 크게 나타남"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch2. 수요와 공급 구조
# ═══════════════════════════════════════════════════════════════════════════
section_slide(2, "수요와 공급 구조")

table_slide(
    "세그먼트별 DRAM 수요 구조",
    ["세그먼트", "비중(추정)", "주요 수요처", "특징"],
    [
        ["서버(AI·클라우드)", "~40%", "하이퍼스케일러", "HBM·DDR5 고성능 수요 급증"],
        ["모바일",            "~30%", "스마트폰 OEM",   "플래그십↑ 평균 탑재량 증가"],
        ["PC·노트북",         "~20%", "OEM·소비자",     "교체 사이클 의존, AI PC 전환"],
        ["기타(차량·산업)",   "~10%", "자동차·가전",    "장기계약 비중 높아 사이클 완충"],
    ],
    note="서버 비중 지속 증가 → 전체 DRAM 수요의 AI 의존도 상승"
)

table_slide(
    "Bit Growth 수급 전망",
    ["연도", "DRAM Bit 수요 성장", "DRAM Bit 공급 성장", "NAND Bit 수요", "NAND Bit 공급"],
    [
        ["2022", "+17%", "+19%", "+25%", "+33%"],
        ["2023", "+14%", "+5%",  "+18%", "+3%"],
        ["2024", "+20%", "+18%", "+22%", "+20%"],
        ["2025E","~22%", "~18%", "~25%", "~22%"],
        ["2026E","~20%", "~16%", "~22%", "~20%"],
    ],
    note="2023 공급 과잉 → 2024 회복 → 2025~26 다시 수요 > 공급 국면 기대"
)

content_slide(
    "3사 과점 구조와 경제적 해자",
    [
        (0, "DRAM은 삼성전자·SK하이닉스·마이크론 3사 과점 (합산 점유율 ~95%)"),
        (1, "진입장벽: 초기 투자 15~20조 원 이상, 공정 기술 노하우 수십 년 필요"),
        (1, "중국 CXMT 등 후발 업체는 DDR4 수준에 머뭄 (2024 기준)"),
        (0, "NAND는 6사 경쟁 구도 (삼성·SK·마이크론·키옥시아·WD·양쯔메모리)"),
        (1, "3D NAND 200단 이상 기술 개발비 부담으로 인수합병 논의 지속"),
        (0, "과점 = 가격 사이클 관리 능력"),
        (1, "3사 모두 감산 동조 시 공급 충격 증폭 → 가격 급등 패턴"),
        (1, "삼성이 'Chicken Game' 돌입 시 경쟁사 수익성 급격히 악화"),
    ],
    note="DRAM 과점은 사이클 하락기 단기 고통, 상승기 폭발적 수익 구조 형성"
)

content_slide(
    "HBM — 이중 수급 구조의 등장",
    [
        (0, "High Bandwidth Memory (HBM) = AI 가속기 필수 부품"),
        (1, "엔비디아 H100/H200/B200 GPU에 HBM3·HBM3e 탑재"),
        (1, "기존 DRAM과 별도 수급 사이클 — 할당량(Allocation) 계약 방식"),
        (0, "HBM이 기존 DRAM 시장에 미치는 영향"),
        (1, "HBM 1GB = DDR5 약 3~4GB 생산 여력 소모 → 범용 DRAM 공급 감소"),
        (1, "SK하이닉스 HBM 점유율 ~50% (2024) → 범용 DRAM ASP 프리미엄"),
        (0, "리스크"),
        (1, "엔비디아 공급 집중 → 단일 고객 의존 리스크"),
        (1, "HBM 수요 둔화 시 전환 비용과 재고 조정 동시 발생 가능"),
    ],
    note="HBM = 메모리 시장의 새로운 프리미엄 세그먼트 — 수급 구조 이원화"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch3. 가격 메커니즘
# ═══════════════════════════════════════════════════════════════════════════
section_slide(3, "가격 메커니즘")

table_slide(
    "Spot 가격 vs Contract 가격",
    ["구분", "Spot (현물)", "Contract (고정계약)"],
    [
        ["계약 주기", "즉시 거래",            "분기/반기 단위"],
        ["가격 결정", "수급 즉각 반영",       "분기초 협상 결정"],
        ["변동성",   "매우 높음",             "상대적으로 낮음"],
        ["주요 거래자","소규모 유통·브로커",  "대형 고객사(서버OEM 등)"],
        ["선행성",   "Contract 선행 (2~4주)", "후행 지표"],
        ["활용",     "시장 방향성 빠른 파악", "실제 수익 예측"],
    ],
    note="Spot이 먼저 움직인다 → Spot↑ 지속 시 Contract 인상 협상 시작"
)

content_slide(
    "스프레드 전략 — 원가와 가격의 격차",
    [
        (0, "메모리 수익성 = 판매가(ASP) - 원가(Cost)"),
        (1, "원가는 공정 세대(nm 또는 레이어)에 따라 구조적으로 하락"),
        (1, "가격이 원가보다 빠르게 떨어지면 '적자 구조' 진입"),
        (0, "Burn Margin (소각 마진) 개념"),
        (1, "원가보다 낮은 가격에 판매 → 출혈 경쟁"),
        (1, "Burn Margin 깊을수록 감산 압력 증가 → 공급 절벽"),
        (0, "사이클 투자 시사점"),
        (1, "원가 이하 판매 = 사이클 바닥 신호"),
        (1, "원가 구조 가장 우수한 기업(삼성·SK)이 사이클 하락기 생존 유리"),
    ],
    note="Burn Margin 심화 → 감산 → 가격 반등의 사이클 작동"
)

table_slide(
    "Burn Margin 구간별 시장 반응",
    ["Burn Margin 수준", "시장 반응", "투자 시사점"],
    [
        ["ASP > Cost + 20%", "고수익 구간, 증설 가속",     "사이클 고점 주의"],
        ["ASP ≈ Cost",       "손익분기, 감산 논의 시작",    "중립, 방향성 관찰"],
        ["ASP < Cost 10%",   "출혈 경쟁, 감산 압력 증가",  "바닥 탐색, 분할 매수 고려"],
        ["ASP < Cost 20%+",  "대규모 감산 불가피",         "사이클 반등 임박 신호"],
    ],
    note="Burn Margin이 깊어질수록 공급 자정 메커니즘 작동 → 반등 속도 빠름"
)

content_slide(
    "LTA (장기 공급 계약)의 딜레마",
    [
        (0, "LTA = Long Term Agreement — 가격·물량 사전 확정 계약"),
        (1, "고객: 공급 확보 + 가격 안정 / 제조사: 수요 안정 + 마진 예측"),
        (0, "딜레마 1 — 고점에서 LTA 체결 시"),
        (1, "고객이 시장가 하락기에도 높은 계약가 지불 → 고객 불만"),
        (1, "제조사는 좋지만 다음 계약 협상에서 불리"),
        (0, "딜레마 2 — 저점에서 LTA 체결 시"),
        (1, "공급 부족기에 제조사가 낮은 계약가로 묶임 → 기회 손실"),
        (0, "AI 시대 LTA 변화"),
        (1, "HBM은 사실상 전량 LTA — 할당 계약이 일반화"),
        (1, "범용 DRAM은 LTA 비중 감소 → 스팟 노출 증가"),
    ],
    note="LTA는 양날의 검 — 시장 사이클 파악 후 체결 시점이 관건"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch4. 사이클의 역사
# ═══════════════════════════════════════════════════════════════════════════
section_slide(4, "사이클의 역사")

table_slide(
    "메모리 슈퍼사이클 연표",
    ["시기", "사이클", "핵심 드라이버", "고점 DRAM 가격 변화"],
    [
        ["2010~2013", "PC/서버 업사이클", "PC 보급 확대, 태블릿 등장", "+150%"],
        ["2014~2016", "다운사이클",       "공급 과잉, PC 수요 정체",   "-60%"],
        ["2017~2018", "슈퍼사이클",       "서버 투자 급증, 모바일 고용량", "+200%+"],
        ["2019",      "단기 조정",        "미중 무역분쟁, 재고 소화",   "-40%"],
        ["2020~2021", "코로나 업사이클",  "재택근무·PC·서버 수요",     "+80%"],
        ["2022~2023", "역대 최악 다운",   "금리 인상·소비 위축·재고 급증","-70%+"],
        ["2024~",     "AI 업사이클",      "LLM·HPC·HBM 수요 폭증",    "반등 중"],
    ],
    note="사이클 평균 주기 2~3년, 최근 AI 변수로 패턴 변화 가능성"
)

content_slide(
    "2022~2023 역대 최대 다운사이클",
    [
        (0, "원인 — 3중 악재 동시 발생"),
        (1, "① 코로나 특수 소멸 → PC·스마트폰 수요 급랭"),
        (1, "② 공급망 정상화 → 고객사 재고 급증 (재고 조정 사이클)"),
        (1, "③ 미 연준 공격적 금리 인상 → 데이터센터 투자 위축"),
        (0, "결과"),
        (1, "DRAM Spot 가격 -70%, NAND Spot -65% (2022 고점 대비)"),
        (1, "삼성전자 반도체 부문 2023년 연간 영업적자 약 14조 원"),
        (1, "SK하이닉스 2023 상반기 연속 대규모 적자"),
        (0, "반등 트리거"),
        (1, "2023년 하반기 대규모 감산 + AI 수요 본격화"),
    ],
    note="2022~23 다운사이클은 수요·재고·매크로 3중 악재 — 역사적 학습점"
)

content_slide(
    "AI 업사이클 — 2024년~",
    [
        (0, "ChatGPT 이후 LLM 추론·학습 인프라 투자 폭발"),
        (1, "엔비디아 H100/H200 → GPU당 HBM 탑재량 증가"),
        (1, "하이퍼스케일러(AWS·Azure·GCP) 서버 투자 사이클 재개"),
        (0, "HBM3/HBM3e 수요 급증으로 DRAM 공급 빠듯"),
        (1, "SK하이닉스 HBM 점유율 ~50%, 삼성 ~35%, 마이크론 ~15%"),
        (1, "HBM 생산 여력이 범용 DRAM 공급 감소로 연결"),
        (0, "업사이클 지속 조건"),
        (1, "AI 투자 지속 + 소비자 IT 수요 회복 (PC·모바일)"),
        (1, "중국 CXMT 등 후발 공급 확대 속도가 수요 증가 속도 내 유지"),
    ],
    note="AI 업사이클 = HBM 프리미엄 + 범용 DRAM 공급 축소 이중 효과"
)

content_slide(
    "변곡점 인식의 비대칭성",
    [
        (0, "사이클 바닥보다 고점이 먼저 보인다는 통념은 틀림"),
        (1, "실제로 고점·바닥 모두 사후에 명확해짐"),
        (1, "선행지표: 재고 weeks, Spot-Contract 스프레드, Burn Margin"),
        (0, "바닥 신호"),
        (1, "Burn Margin 심화 + 주요 업체 감산 발표"),
        (1, "재고 Weeks 정점 통과 + 고객사 긴급 발주 시작"),
        (0, "고점 신호"),
        (1, "신규 Capa 투자 러시 + LTA 계약 급증"),
        (1, "Spot-Contract 프리미엄 축소 + 중국 공급 증가"),
        (0, "투자 시사점"),
        (1, "변곡점 인식 후 3~6개월 안에 포지션 진입이 최적"),
    ],
    note="선행지표 조합으로 변곡점 확률적 파악 → 비대칭 수익 기회"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch5. 미래 전망
# ═══════════════════════════════════════════════════════════════════════════
section_slide(5, "미래 전망")

table_slide(
    "Commodity → Custom 전환 트렌드",
    ["구분", "Commodity DRAM", "Custom/HBM"],
    [
        ["제품 형태",  "표준 규격 (DDR5 등)",   "고객 맞춤 설계"],
        ["마진",       "낮음 (사이클 종속)",     "높음 (5~10배 ASP 프리미엄)"],
        ["고객 관계",  "경쟁 입찰",              "Joint Development / LTA"],
        ["진입장벽",   "공정 미세화",            "공정 + 패키징(TSV) 기술"],
        ["성장 드라이버", "수량(Bit) 증가",      "ASP + 수량 동시 성장"],
        ["대표 사례",  "DDR5, LPDDR5",           "HBM3e, CXL Memory"],
    ],
    note="HBM·CXL 등 Custom 비중 확대 → 메모리 마진 구조적 개선 가능"
)

content_slide(
    "DDR5 전환과 AI PC",
    [
        (0, "DDR5 전환 가속 — 2024~2026년이 핵심 교체 구간"),
        (1, "인텔 14세대(Meteor Lake), AMD Zen5 → DDR5 기본 채택"),
        (1, "DDR5 가격 프리미엄 축소로 전환 비용 부담 완화"),
        (0, "AI PC = 엣지 AI 추론을 위한 고성능 DRAM 탑재"),
        (1, "온디바이스 LLM 구동 → LPDDR5X 고용량·고대역폭 요구"),
        (1, "NPU 내장 CPU + 고용량 DRAM = AI PC 필수 조건"),
        (0, "전망"),
        (1, "2025~26년 PC 교체 사이클 + AI PC 전환 동시 발생"),
        (1, "DRAM 탑재량: 일반 PC 16GB → AI PC 32GB+ 확대"),
    ],
    note="DDR5 전환 + AI PC = DRAM 평균 탑재량 증가 → Bit 수요 구조적 상승"
)

content_slide(
    "중국 변수 — 위협과 한계",
    [
        (0, "CXMT(창신메모리) — 중국의 DRAM 독자 개발"),
        (1, "현 수준: DDR4 19nm급 (삼성 1z nm vs 격차 2~3세대)"),
        (1, "EUV 장비 수입 차단 → 첨단 공정 전환 제약"),
        (0, "YMTC(양쯔메모리) — 중국 NAND"),
        (1, "232단 3D NAND 자체 개발 (Xtacking 아키텍처)"),
        (1, "미국 제재로 글로벌 판매 제한, 주로 국내 공급"),
        (0, "중국 변수의 시장 영향"),
        (1, "단기: 범용 DRAM·NAND 가격 압박 (중저가 시장 점유)"),
        (1, "중기: HBM·첨단 DRAM 진입 어려움 → 프리미엄 시장 영향 제한"),
        (1, "장기: EUV·패키징 기술 격차 좁혀질 가능성 모니터링 필요"),
    ],
    note="중국 위협 = 범용 시장 가격 압박 + 첨단 시장은 기술 장벽으로 방어 중"
)

table_slide(
    "2026년 메모리 시장 위치 진단",
    ["지표", "현황(2025)", "2026년 전망"],
    [
        ["DRAM Bit 수요 성장", "~22%",           "~20% (AI 지속, PC 회복)"],
        ["DRAM Bit 공급 성장", "~18%",           "~16% (감산 여파 + HBM 전환)"],
        ["수급",               "소폭 타이트",    "타이트 지속 (공급 < 수요)"],
        ["HBM 비중",           "DRAM의 ~15%",    "~20%+ (B200·GB200 채택)"],
        ["중국 공급 영향",     "제한적",         "범용 시장 일부 압박 가능"],
        ["NAND 수급",          "균형~소폭 과잉", "균형 (감산 효과 반영)"],
    ],
    note="2026년은 업사이클 지속 국면 — HBM 수요가 핵심 변수"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch6. 시장 심층 분석
# ═══════════════════════════════════════════════════════════════════════════
section_slide(6, "시장 심층 분석")

table_slide(
    "2026 DRAM 시장 전망 상세",
    ["구분", "2025E", "2026E", "YoY"],
    [
        ["전체 DRAM 매출(B$)", "$95~100", "$110~120", "+15~20%"],
        ["HBM 매출(B$)",       "$20~25",  "$35~40",   "+50~60%"],
        ["범용 DRAM 매출",     "$70~80",  "$75~85",   "+5~10%"],
        ["DDR5 비중",          "~55%",    "~70%",     "+15%p"],
        ["DRAM ASP(평균)",     "회복세",  "소폭 상승", "안정적"],
    ],
    note="HBM 고성장이 전체 DRAM 시장 매출 견인 — 범용 성장은 완만"
)

table_slide(
    "DRAM 사이클 역사표 (확장)",
    ["사이클", "기간", "DRAM 가격 변화", "주요 원인"],
    [
        ["업사이클 1", "2010~2013", "+150%",   "PC 보급, 스마트폰 초기"],
        ["다운 1",     "2014~2016", "-60%",    "공급 과잉, PC 정체"],
        ["슈퍼사이클", "2017~2018", "+200%+",  "서버 투자 폭증, 모바일 고용량화"],
        ["단기 조정",  "2019",      "-40%",    "미중 무역분쟁"],
        ["코로나 업",  "2020~2021", "+80%",    "WFH·PC·서버 특수"],
        ["역대 최악",  "2022~2023", "-70%+",   "3중 악재(수요·재고·금리)"],
        ["AI 업",      "2024~",     "반등 중", "LLM·HBM·서버 투자 재개"],
    ],
    note="평균 사이클 2~3년, AI 변수로 향후 사이클 패턴 변화 가능성"
)

content_slide(
    "Snowball Effect — 감산의 증폭 메커니즘",
    [
        (0, "Snowball Effect = 초기 감산이 공급 부족을 과도하게 증폭시키는 현상"),
        (1, "①감산 발표 → ②고객 긴급 재고 확보 → ③가격 급등"),
        (1, "④추가 감산 선언 → ⑤공급 절벽 → ⑥가격 급등 지속"),
        (0, "Re-entrant Flow와 결합 시 효과 배가"),
        (1, "초기 5~10% 감산이 실제 Bit 공급 15~20% 감소로 이어질 수 있음"),
        (0, "투자 시사점"),
        (1, "초기 감산 신호 포착 시 빠른 포지션 구축이 유리"),
        (1, "감산 발표 → 수개월 후 실제 공급 감소 → 가격 반등 시차 활용"),
    ],
    note="감산 발표 ≠ 즉각 공급 감소 — 시차(3~6개월) 파악이 핵심"
)

content_slide(
    "생산 제약 — 기술·설비·인력",
    [
        (0, "EUV 노광 장비 부족"),
        (1, "ASML EUV 장비 수요 > 공급 → 납기 2~3년 대기"),
        (1, "첨단 DRAM(1β nm 이하)·HBM TSV 공정에 필수"),
        (0, "TSV(Through-Silicon Via) 패키징 제약"),
        (1, "HBM 적층을 위한 TSV 공정 수율이 핵심 병목"),
        (1, "SK하이닉스·삼성 외 TSV 대량생산 가능 업체 부재"),
        (0, "클린룸 확장 vs 기술 전환 딜레마"),
        (1, "신규 팹 건설 비용 15~20조 원, 완공까지 2~3년"),
        (1, "기존 팹을 HBM 라인으로 전환 시 범용 DRAM 생산 축소"),
    ],
    note="생산 제약 3중주: EUV 장비·TSV 수율·팹 전환 비용"
)

table_slide(
    "HBM 세대별 스펙 비교",
    ["세대", "대역폭", "용량", "적층 수", "주요 채택 제품"],
    [
        ["HBM2e", "460 GB/s", "16GB", "8단", "A100"],
        ["HBM3",  "819 GB/s", "24GB", "12단","H100"],
        ["HBM3e", "1.2 TB/s", "36GB", "12단","H200, MI300X"],
        ["HBM4",  "2.0 TB/s+","48GB+","16단","B200, GB200 (예정)"],
    ],
    note="HBM4부터 베이스다이 로직 통합 → SK하이닉스·삼성의 기술 격차 확대"
)

content_slide(
    "NAND Flash 시장 심층 분석",
    [
        (0, "NAND 시장 구조 — 6사 경쟁, 차별화 어려움"),
        (1, "삼성(~32%), SK하이닉스+솔리다임(~20%), 마이크론(~16%), 키옥시아(~15%)"),
        (1, "QLC(4-bit) 확대로 GB당 원가 하락, 엔터프라이즈 SSD 시장 성장"),
        (0, "3D NAND 레이어 경쟁"),
        (1, "현재 선두: 삼성 290단, SK 238단, 마이크론 232단"),
        (1, "200단 이상 적층 시 본딩 기술(Wafer-to-Wafer) 필요 → 진입장벽"),
        (0, "NAND 투자 시사점"),
        (1, "DRAM 대비 사이클 진폭이 크고 회복 속도가 느림"),
        (1, "데이터센터 SSD 수요(AI 학습·추론용 스토리지) 구조적 성장"),
    ],
    note="NAND = 경쟁 심화 + 레이어 경쟁 → 기술 선두 업체 선별 중요"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch7. 투자 아이디어
# ═══════════════════════════════════════════════════════════════════════════
section_slide(7, "투자 아이디어")

content_slide(
    "LTA 구조와 투자 기회",
    [
        (0, "LTA(장기 공급 계약) 체결 시점이 투자 수익 결정"),
        (1, "사이클 바닥 LTA → 제조사 저수익 / 고객 이익"),
        (1, "사이클 고점 LTA → 제조사 수익 극대화 / 고객 리스크"),
        (0, "HBM LTA의 특수성"),
        (1, "엔비디아·AMD·인텔과 연간~복수연도 할당 계약"),
        (1, "물량 고정이지만 ASP는 협상 — 기술 로드맵 리스크 공유"),
        (0, "투자 전략"),
        (1, "LTA 체결 뉴스 모니터링 → 공급자 유리 조건인지 확인"),
        (1, "사이클 하락기 LTA 체결 후 업사이클 ASP 인상 여력 파악"),
    ],
    note="LTA 조건 = 사이클 위치 바로미터 — 투자 선행지표로 활용"
)

table_slide(
    "Dual Market — HBM vs 범용 DRAM",
    ["구분", "HBM", "범용 DRAM (DDR5 등)"],
    [
        ["수요 드라이버",  "AI 가속기 (엔비디아 등)", "서버·PC·모바일"],
        ["가격 결정",      "협상 (LTA·할당)",         "Spot·Contract"],
        ["ASP 프리미엄",   "DDR5 대비 5~10배",        "기준"],
        ["마진",           "매우 높음",                "사이클 종속"],
        ["수요 가시성",    "높음 (2~3년 로드맵)",      "분기별 변동"],
        ["리스크",         "단일 고객 의존",           "Commodity 가격 변동"],
    ],
    note="포트폴리오 관점: HBM 高마진 안정 + 범용 DRAM 사이클 업사이드 조합"
)

two_col_slide(
    "AI Scaling Law · Memory Wall · KV Cache",
    [
        (0, "AI Scaling Law"),
        (1, "모델 파라미터↑ → 성능↑"),
        (1, "GPT-3 175B → GPT-4 ~1T 파라미터"),
        (1, "학습·추론 메모리 요구량 기하급수적 증가"),
        (0, "Memory Wall"),
        (1, "CPU/GPU 연산 속도 > 메모리 대역폭"),
        (1, "HBM이 부분 해결책 — 대역폭 1TB/s+"),
        (1, "CXL Memory로 용량 문제 보완"),
    ],
    [
        (0, "KV Cache 메모리 수요"),
        (1, "LLM 추론 시 KV Cache = 문맥 저장소"),
        (1, "컨텍스트 길이 2배 → KV Cache 2배"),
        (1, "100K 토큰 컨텍스트 → 수십 GB DRAM 소모"),
        (0, "투자 시사점"),
        (1, "LLM 추론 서버 = DRAM 대용량화 수혜"),
        (1, "고대역폭(HBM) + 고용량(CXL) 동시 필요"),
        (1, "메모리 기업 구조적 성장 모멘텀"),
    ],
    note="AI = 메모리 수요의 질적 변화 — 대역폭 + 용량 동시 프리미엄"
)

content_slide(
    "Revaluation — 메모리의 가치 재평가",
    [
        (0, "전통적 메모리 밸류에이션: PBR 0.5~2배 (사이클 종속)"),
        (0, "AI 시대 재평가 논리"),
        (1, "HBM = AI 인프라 필수 부품 → 성장주 특성 부여"),
        (1, "범용 DRAM도 AI·엣지 수요로 탑재량 구조적 증가"),
        (1, "Custom 메모리 비중 확대 → 마진 안정성 향상"),
        (0, "밸류에이션 변화 방향"),
        (1, "SK하이닉스: HBM 점유율 프리미엄 → PBR 2~3배 정당화 논의"),
        (1, "마이크론: HBM 후발주자지만 성장 스토리 유효"),
        (1, "삼성: HBM 경쟁력 회복 여부가 밸류 디스카운트 해소 관건"),
    ],
    note="HBM 점유율 = 메모리 기업 밸류에이션의 새로운 핵심 지표"
)

two_col_slide(
    "삼성전자 vs SK하이닉스 — 투자 비교",
    [
        (0, "삼성전자"),
        (1, "DRAM·NAND·파운드리 종합"),
        (1, "HBM 수율 이슈로 점유율 하락"),
        (1, "2024 HBM3e 납품 지연 논란"),
        (1, "엔비디아 인증 여부 핵심 모니터링"),
        (0, "밸류에이션"),
        (1, "PBR 1배 이하 (역사적 저점)"),
        (1, "HBM 회복 시 디스카운트 해소"),
        (1, "파운드리 적자 지속이 불안 요인"),
    ],
    [
        (0, "SK하이닉스"),
        (1, "HBM 글로벌 점유율 ~50%"),
        (1, "엔비디아 독점 공급자 지위"),
        (1, "HBM3e → HBM4 로드맵 선도"),
        (1, "범용 DRAM도 수익성 회복"),
        (0, "밸류에이션"),
        (1, "PBR 1.5~2배 (HBM 프리미엄)"),
        (1, "AI 수요 지속 시 추가 리레이팅"),
        (1, "단일 고객(엔비디아) 의존 리스크"),
    ],
    note="삼성: HBM 회복 옵션 매력 / SK: HBM 모멘텀 지속 베팅"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch8. 핵심 프레임 5가지
# ═══════════════════════════════════════════════════════════════════════════
section_slide(8, "핵심 프레임 5가지")

content_slide(
    "프레임 1 — 사이클 위치 파악",
    [
        (0, "메모리 투자의 알파는 사이클 위치 파악에서 나온다"),
        (1, "재고 Weeks (업계 평균): 4~6주 = 정상 / 8주+ = 과잉 / 2주↓ = 부족"),
        (1, "Spot-Contract 스프레드: 확대 = 공급 부족 / 축소 = 완화"),
        (1, "Burn Margin: 적자 판매 = 바닥 근접 신호"),
        (0, "사이클 위치 → 투자 전략 매핑"),
        (1, "바닥 탐색기: 감산 발표 + Burn Margin 심화 → 분할 매수"),
        (1, "업사이클 초입: 재고 소진 + 가격 반등 확인 → 추가 매수"),
        (1, "고점 주의: 신규 Capa 투자 러시 + 스프레드 축소 → 비중 축소"),
    ],
    note="프레임 1: 재고+스프레드+Burn Margin 3지표 조합으로 사이클 위치 판단"
)

content_slide(
    "프레임 2 — 기술 세대 전환 속도",
    [
        (0, "공정 세대가 바뀔 때 수급이 재편된다"),
        (1, "DDR4→DDR5 전환: 초기 공급 부족 → 수율 안정화 → 공급 정상화"),
        (1, "HBM2e→HBM3→HBM3e→HBM4: 각 세대 전환 시 공급 병목"),
        (0, "기술 전환 속도 모니터링"),
        (1, "수율 개선 속도: 빠를수록 공급 조기 정상화 → 가격 하락 압력"),
        (1, "레거시 vs 첨단 라인 전환 속도가 공급 조절 핵심"),
        (0, "투자 시사점"),
        (1, "기술 선도 기업이 세대 전환기 초기에 독점 공급 → 프리미엄 마진"),
        (1, "후발 기업의 추격 속도가 프리미엄 마진 유지 기간 결정"),
    ],
    note="프레임 2: 기술 세대 전환 = 공급 병목 → 선도 기업 프리미엄 기회"
)

content_slide(
    "프레임 3 — 수요의 질적 변화",
    [
        (0, "수요는 양(Bit)과 질(ASP·마진)로 나눠 봐야 한다"),
        (1, "과거: Bit 수요 증가 = 단가 하락 + 물량 증가 (마진 불변)"),
        (1, "현재: HBM·CXL → ASP 프리미엄으로 Bit 성장 이상의 매출 성장"),
        (0, "AI가 바꾼 수요 구조"),
        (1, "추론 서버: 고대역폭(HBM) + 고용량(LPDDR5X·CXL)"),
        (1, "엣지 AI: 모바일·PC DRAM 용량 구조적 증가"),
        (0, "수요 질적 변화 체크포인트"),
        (1, "HBM 매출 비중 추이 (15%→20%+)"),
        (1, "DDR5·LPDDR5X 비중 전환 속도"),
        (1, "서버 DRAM 평균 탑재량 GB 추이"),
    ],
    note="프레임 3: Bit 성장 + ASP 상승 = 매출·마진 동시 성장 구조 확인"
)

content_slide(
    "프레임 4 — 공급 제약의 지속성",
    [
        (0, "공급 제약은 일시적인가, 구조적인가?"),
        (1, "일시적 제약: 수율 이슈, 단기 장비 납기 지연"),
        (1, "구조적 제약: EUV 장비 부족, TSV 공정 기술 장벽"),
        (0, "현재 공급 제약의 성격 진단"),
        (1, "HBM TSV 수율 제약 → 구조적 (2~3년 해소 전망)"),
        (1, "EUV 장비 부족 → 구조적 (ASML 생산 한계)"),
        (1, "중국 수출 규제 → 정책 변수 (불확실성 높음)"),
        (0, "투자 시사점"),
        (1, "구조적 공급 제약 = 사이클 고점이 더 높고 더 오래 지속"),
        (1, "일시적 제약 = 빠른 해소 후 공급 정상화 → 단기 프리미엄만"),
    ],
    note="프레임 4: 공급 제약의 지속성 파악이 투자 기간 설정에 핵심"
)

content_slide(
    "프레임 5 — 거시 환경과 메모리 사이클",
    [
        (0, "금리·환율·무역 환경이 메모리 사이클을 증폭시킨다"),
        (1, "금리 인상기: 데이터센터 투자 위축 → DRAM 수요 감소"),
        (1, "금리 인하기: IT 투자 재개 + 소비자 IT 수요 회복"),
        (0, "환율 영향"),
        (1, "원/달러 약세: 한국 메모리 기업 달러 매출 원화 환산 증가"),
        (1, "엔/달러 약세: 일본 키옥시아 원가 경쟁력 약화"),
        (0, "무역 분쟁·수출 규제"),
        (1, "미중 갈등 → 중국향 메모리 제품 제한 → 단기 공급 왜곡"),
        (1, "반도체법·보조금 → 지역 생산 기지 분산 → 장기 구조 변화"),
        (0, "투자 시사점"),
        (1, "거시 = 메모리 사이클의 외부 증폭기 — 방향은 수급이 결정"),
    ],
    note="프레임 5: 거시는 증폭기, 수급이 방향 — 거시 논리에 매몰 금지"
)

# ═══════════════════════════════════════════════════════════════════════════
# Ch9. Risk 분석
# ═══════════════════════════════════════════════════════════════════════════
section_slide(9, "Risk 분석")

content_slide(
    "AI 수익성 리스크 — 투자 대비 수익 불확실성",
    [
        (0, "하이퍼스케일러의 AI 인프라 투자 ROI 불확실"),
        (1, "AWS·Azure·GCP 합산 2024 AI Capex ~ $200B 예상"),
        (1, "수익화 모델(AI 서비스 매출) 아직 초기 단계"),
        (0, "AI 버블 논란"),
        (1, "닷컴 버블 유사성: 과잉 투자 → 조정 → 생산성 혁명 순서로 전개"),
        (1, "단기 AI 투자 둔화 시 HBM·서버 DRAM 수요 급감 가능"),
        (0, "리스크 관리 포인트"),
        (1, "하이퍼스케일러 Capex 가이던스 분기별 모니터링"),
        (1, "엔비디아 GPU 주문 취소·연기 징후 조기 감지"),
        (1, "AI 서비스 매출화 속도가 투자 지속성 결정"),
    ],
    note="AI 수익성 리스크 = HBM 수요의 최대 꼬리 위험"
)

content_slide(
    "빅테크 자체 칩 개발 리스크",
    [
        (0, "빅테크의 AI 칩 내재화 가속"),
        (1, "구글 TPU, 아마존 Trainium, 메타 MTIA, 애플 M·A 시리즈"),
        (1, "자체 칩 확대 → 엔비디아 GPU 의존도 축소 가능"),
        (0, "메모리에 미치는 영향"),
        (1, "자체 칩도 HBM/LPDDR 필요 → 메모리 수요 자체는 유지"),
        (1, "단, 엔비디아 중심 HBM 수요 집중 구조 변화 가능"),
        (1, "HBM 공급업체 다변화 압력 발생 가능"),
        (0, "대응 전략"),
        (1, "자체 칩 업체별 메모리 스펙·공급 계약 동향 추적"),
        (1, "SK하이닉스 외 마이크론·삼성의 커스텀 HBM 협상 현황"),
    ],
    note="빅테크 칩 내재화 = 엔비디아 리스크이나 메모리 수요 자체는 방어적"
)

content_slide(
    "FAB 증설 과잉 리스크",
    [
        (0, "업사이클 → 공격적 증설 → 과잉 공급 → 다운사이클 반복"),
        (1, "2017~18 슈퍼사이클 후 2019 조정의 교훈"),
        (1, "2020~21 코로나 특수 → 2022~23 역대 최악 다운사이클"),
        (0, "현재 증설 현황 (2024~2026)"),
        (1, "삼성: 평택 P4·P5 라인 HBM 전환 중심"),
        (1, "SK하이닉스: 청주·용인 첨단 팹 투자"),
        (1, "마이크론: 미국 아이다호·뉴욕 Fab 보조금 투자"),
        (0, "과잉 증설 방지 요인"),
        (1, "HBM 전환 투자 = 범용 DRAM 감산 효과 → 자동 조절"),
        (1, "EUV 장비 납기 제약 → 급격한 Capa 확장 제한"),
    ],
    note="FAB 증설 리스크는 있으나 HBM 전환·EUV 제약이 완충 역할"
)

table_slide(
    "공급 증가 vs 수요 증가 — 종합 결론",
    ["시나리오", "공급 성장", "수요 성장", "수급", "투자 판단"],
    [
        ["Base (가능성 50%)",  "~16%", "~20%", "타이트",    "긍정적 — 업사이클 지속"],
        ["Bull (가능성 25%)",  "~14%", "~24%", "매우 타이트","매우 긍정적 — HBM·가격 폭등"],
        ["Bear (가능성 25%)",  "~20%", "~15%", "과잉",       "부정적 — AI 투자 급감 시나리오"],
    ],
    note="Base 시나리오 기준 2026년 메모리 업사이클 지속 — AI 수요 모니터링이 핵심"
)

content_slide(
    "종합 Risk Map",
    [
        (0, "고영향·고확률 리스크"),
        (1, "① 하이퍼스케일러 AI Capex 둔화 → HBM 수요 급감"),
        (1, "② 중국 CXMT 범용 DRAM 공급 급증"),
        (0, "고영향·저확률 리스크"),
        (1, "③ 엔비디아 GPU 수요 붕괴 (경쟁사 급부상)"),
        (1, "④ 글로벌 경기 침체 → IT 지출 전반 위축"),
        (0, "저영향·고확률 리스크"),
        (1, "⑤ NAND 가격 변동성 (범용 경쟁 지속)"),
        (1, "⑥ 원/달러 환율 불리한 방향 전환"),
        (0, "투자 결론"),
        (1, "Risk 대비 Return: Base 시나리오 기준 긍정적"),
        (1, "핵심 모니터링: 하이퍼스케일러 Capex + HBM 수요 신호"),
    ],
    note="리스크는 인식하되 과잉 반응 금물 — 사이클 기반 분할 대응이 최적"
)

# ── 마지막 요약 슬라이드 ───────────────────────────────────────────────────
content_slide(
    "핵심 요약 — 메모리 산업 투자 프레임워크",
    [
        (0, "① 사이클 파악: 재고·스프레드·Burn Margin 3지표 조합"),
        (0, "② 기술 전환: 세대 전환 초기가 선도 기업 프리미엄 극대화 구간"),
        (0, "③ 수요 질적 변화: Bit + ASP 동시 상승 = 메모리 리레이팅"),
        (0, "④ 공급 제약: EUV·TSV 구조적 병목 → 사이클 고점 연장"),
        (0, "⑤ 거시 증폭: 금리·무역은 방향이 아닌 증폭기"),
        (0, "2026년 포지션"),
        (1, "업사이클 지속 국면 — HBM 선도 기업 비중 확대"),
        (1, "삼성 HBM 회복 시 추가 옵션 가치"),
        (1, "AI Capex 모니터링 — 핵심 선행지표 주간 점검 권고"),
    ],
    note="메모리 = 사이클 + AI 구조 변화의 교차점 — 프레임워크 기반 투자"
)

# ── 저장 ──────────────────────────────────────────────────────────────────
OUT = r"c:\Users\yh900\SemiCon\memory_industry_v2.pptx"
prs.save(OUT)
print(f"저장 완료: {OUT}")
print(f"슬라이드 수: {len(prs.slides)}")
